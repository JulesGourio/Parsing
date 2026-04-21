import io
import json
import os
import re
import shutil
import statistics
import subprocess
import tempfile
import time
import unicodedata
from difflib import SequenceMatcher
from importlib import metadata as importlib_metadata
from pathlib import Path
from typing import Any, Callable, Dict, Iterable, List, Optional

import pandas as pd
import pdfplumber
from docx import Document
from docx.oxml.ns import qn
from lxml import etree
from pptx import Presentation

try:
    import numpy as np
except Exception:  # pragma: no cover - optional dependency
    np = None

try:
    from PIL import Image, ImageOps
except Exception:  # pragma: no cover - optional dependency
    Image = None
    ImageOps = None

try:
    import pypdfium2 as pdfium
except Exception:  # pragma: no cover - optional dependency
    pdfium = None

try:
    import pytesseract
except Exception:  # pragma: no cover - optional dependency
    pytesseract = None

try:
    from rapidocr_onnxruntime import RapidOCR
except Exception:  # pragma: no cover - optional dependency
    RapidOCR = None

from .constants import (
    ANTIWORD_BIN,
    ANTIWORD_CPU_SECONDS,
    ANTIWORD_MEMORY_MB,
    ANTIWORD_SHARE_DIR,
    EXCEL_INCLUDE_ROW_CONTEXT,
    EXCEL_TABLE_PREVIEW_ROWS,
    LOCAL_PARSE_MAX_RETRIES,
    OCR_ENABLED,
    OCR_ENGINE_PRIORITY,
    OCR_IMAGE_PREPROCESS,
    OCR_MAX_IMAGE_SIDE_PX,
    OCR_MAX_PDF_OCR_ATTEMPTS,
    OCR_MAX_PDF_PAGES,
    OCR_MIN_CONFIDENCE,
    OCR_MIN_IMAGE_SIDE_PX,
    OCR_MIN_PDF_PAGE_TEXT_CHARS,
    OCR_PDF_OCR_TIME_BUDGET_RATIO,
    OCR_PDF_FALLBACK_ENABLED,
    OCR_PDF_RENDER_DPI,
    OCR_SUPPLEMENT_LINE_SIMILARITY,
    OCR_SUPPLEMENT_TEXT_SIMILARITY,
    OCR_TABLE_MIN_BOXES,
    OCR_TABLE_MIN_COLUMNS,
    OCR_TABLE_MIN_DENSITY,
    OCR_TABLE_RECONSTRUCTION_ENABLED,
    OCR_DIAGRAM_SUMMARY_ENABLED,
    OCR_DIAGRAM_MIN_BOXES,
    OCR_TESSERACT_BIN,
    OCR_TESSERACT_LANG,
    OCR_TESSERACT_OEM,
    OCR_TESSERACT_PSM,
    OCR_TIMEOUT_SECONDS,
    PDF_EXCLUDE_TABLE_TEXT,
    MAX_CHUNK_CHARS,
    MAX_CSV_ROWS,
    MAX_CSV_TABLE_PREVIEW_ROWS,
    MAX_DOCUMENT_BYTES,
    MAX_ERROR_MESSAGE_CHARS,
    MAX_EXCEL_ROWS_PER_SHEET,
    MAX_EXCEL_SHEETS,
    MAX_FALLBACK_TEXT_CHARS,
    MAX_PDF_PAGES,
    MAX_PPTX_SLIDES,
    PARSE_TIMEOUT_SECONDS,
    TEXTLIKE_EXTENSIONS,
)
from .retry import is_retryable_error
from .text_utils import normalize_text, rows_to_markdown_table, safe_decode, strip_html_tags

_LIB_VERSION_CACHE: Dict[str, str] = {}
_BINARY_SIGNATURE_FORMATS = {"pdf", "jpeg", "png", "tiff", "webp", "ole2", "ooxml", "rtf"}
_IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "tif", "tiff", "bmp", "webp"}
_RAPID_OCR_ENGINE: Optional[Any] = None


def _response(text: str, parser_error: Optional[str], parser_strategy: str) -> Dict[str, Any]:
    return {
        "text": text,
        "parser_error": parser_error,
        "parser_strategy": parser_strategy,
        "ocr_attempted": False,
        "ocr_used": False,
        "ocr_engine_trace": None,
        "ocr_pages": 0,
        "ocr_supplement_pages": 0,
    }


def _with_ocr_metadata(
    result: Dict[str, Any],
    *,
    attempted: bool,
    used: bool,
    trace: Optional[str],
    pages: int,
    supplement_pages: int,
) -> Dict[str, Any]:
    result["ocr_attempted"] = bool(attempted)
    result["ocr_used"] = bool(used)
    result["ocr_engine_trace"] = normalize_text(trace or "") or None
    result["ocr_pages"] = max(0, _safe_int(pages))
    result["ocr_supplement_pages"] = max(0, _safe_int(supplement_pages))
    return result


def _sanitize_error_message(message: str) -> str:
    if not message:
        return ""

    cleaned = str(message).replace("\r", " ").replace("\n", " ").strip()
    cleaned = re.sub(r"[A-Za-z]:\\[^\s]+", "<path>", cleaned)
    cleaned = re.sub(r"/(?:[^/\s]+/)*[^/\s]+", "<path>", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned)

    if len(cleaned) > MAX_ERROR_MESSAGE_CHARS:
        cleaned = f"{cleaned[: MAX_ERROR_MESSAGE_CHARS - 3]}..."
    return cleaned


def _error(code: str, message: str) -> str:
    return f"{code}: {_sanitize_error_message(message)}"


def _is_timed_out(started_at: float, timeout_seconds: float = PARSE_TIMEOUT_SECONDS) -> bool:
    return (time.time() - started_at) > timeout_seconds


def _normalize_extension(extension: str) -> str:
    ext = (extension or "").lower().strip()
    return ext[1:] if ext.startswith(".") else ext


def _detect_binary_format(content: bytes) -> str:
    if not content or len(content) < 5:
        return "unknown"
    if content[:5] == b"%PDF-":
        return "pdf"
    if content[:3] == b"\xff\xd8\xff":
        return "jpeg"
    if content[:4] == b"\x89PNG":
        return "png"
    if content[:4] in {b"II*\x00", b"MM\x00*"}:
        return "tiff"
    if content[:4] == b"RIFF" and content[8:12] == b"WEBP":
        return "webp"
    if content[:4] == b"\xd0\xcf\x11\xe0":
        return "ole2"
    if content[:4] == b"PK\x03\x04":
        return "ooxml"
    if content.lstrip()[:5] == b"{\\rtf":
        return "rtf"
    return "unknown"


def _is_mime_mismatch_for_textlike(extension: str, detected_format: str) -> bool:
    ext = _normalize_extension(extension)
    if ext not in TEXTLIKE_EXTENSIONS:
        return False
    return detected_format in (_BINARY_SIGNATURE_FORMATS - {"rtf"}) or (detected_format == "rtf" and ext != "rtf")


def _truncate_text(text: str, max_chars: int = MAX_FALLBACK_TEXT_CHARS) -> str:
    if len(text) <= max_chars:
        return text
    return f"{text[:max_chars]}\n\n[TRUNCATED at {max_chars} chars]"


def _lib_version(lib_name: str) -> str:
    if lib_name in _LIB_VERSION_CACHE:
        return _LIB_VERSION_CACHE[lib_name]

    try:
        _LIB_VERSION_CACHE[lib_name] = importlib_metadata.version(lib_name)
    except Exception:
        _LIB_VERSION_CACHE[lib_name] = "unknown"
    return _LIB_VERSION_CACHE[lib_name]


def _with_lib_versions(strategy: str, *libs: str) -> str:
    if not libs:
        return strategy
    versions = [f"{lib}={_lib_version(lib)}" for lib in libs]
    return f"{strategy}|libs:{','.join(versions)}"


def _strip_rtf_control_words(decoded: str) -> str:
    # Keep plain textual payload while dropping most RTF control markup.
    no_hex = re.sub(r"\\'[0-9a-fA-F]{2}", " ", decoded)
    no_controls = re.sub(r"\\[a-zA-Z]+-?\d*\s?", " ", no_hex)
    no_groups = no_controls.replace("{", " ").replace("}", " ")
    return normalize_text(no_groups)


def _extract_docx_comments(document: Document) -> List[str]:
    comments: List[str] = []
    comments_part = getattr(document.part, "comments_part", None)
    if comments_part is None:
        return comments

    root = getattr(comments_part, "element", None)
    if root is None:
        return comments

    try:
        for comment in root.xpath(".//*[local-name()='comment']"):
            author = normalize_text(comment.get(qn("w:author"), ""))
            text_fragments = [normalize_text(node.text or "") for node in comment.xpath(".//*[local-name()='t']")]
            text = normalize_text(" ".join([fragment for fragment in text_fragments if fragment]))
            if text:
                label = f"[{author}] " if author else ""
                comments.append(f"{label}{text}")
    except Exception:
        return []

    return comments


def _extract_docx_textboxes(document: Document) -> List[str]:
    seen = set()
    extracted: List[str] = []
    for node in document.element.body.xpath(".//*[local-name()='txbxContent']//*[local-name()='t']"):
        text = normalize_text(node.text or "")
        if text and text not in seen:
            seen.add(text)
            extracted.append(text)
    return extracted


def _collect_docx_image_blobs(document: Document) -> List[bytes]:
    blobs: List[bytes] = []
    seen = set()

    def _collect_from_part(part: Any) -> None:
        rels = getattr(part, "rels", {})
        for rel in rels.values():
            reltype = str(getattr(rel, "reltype", ""))
            if "image" not in reltype:
                continue

            target_part = getattr(rel, "target_part", None)
            blob = getattr(target_part, "blob", None)
            if not blob:
                continue

            signature = (len(blob), blob[:32])
            if signature in seen:
                continue

            seen.add(signature)
            blobs.append(blob)

    _collect_from_part(document.part)
    for section in document.sections:
        try:
            _collect_from_part(section.header.part)
        except Exception:
            pass
        try:
            _collect_from_part(section.footer.part)
        except Exception:
            pass

    return blobs


def _extract_paragraph_with_hyperlinks(paragraph) -> str:
    rels = paragraph.part.rels
    pieces: List[str] = []

    for child in paragraph._p.iterchildren():
        if child.tag == qn("w:r"):
            run_text = "".join([node.text for node in child.xpath(".//*[local-name()='t']") if node.text])
            normalized = normalize_text(run_text)
            if normalized:
                pieces.append(normalized)
        elif child.tag == qn("w:hyperlink"):
            link_text = "".join([node.text for node in child.xpath(".//*[local-name()='t']") if node.text])
            rid = child.get(qn("r:id"))
            href = ""
            if rid and rid in rels:
                href = normalize_text(str(getattr(rels[rid], "target_ref", "")))

            normalized_link_text = normalize_text(link_text)
            if normalized_link_text and href:
                pieces.append(f"{normalized_link_text} ({href})")
            elif normalized_link_text:
                pieces.append(normalized_link_text)
            elif href:
                pieces.append(href)

    text = normalize_text(" ".join(pieces))
    return text or normalize_text(paragraph.text)


def _antiword_preexec_limits() -> Optional[Callable[[], None]]:
    if os.name != "posix":
        return None

    try:
        import resource
    except ImportError:
        return None

    def _apply_limits() -> None:
        try:
            resource.setrlimit(resource.RLIMIT_CPU, (ANTIWORD_CPU_SECONDS, ANTIWORD_CPU_SECONDS))
        except Exception:
            pass

        try:
            max_memory_bytes = ANTIWORD_MEMORY_MB * 1024 * 1024
            resource.setrlimit(resource.RLIMIT_AS, (max_memory_bytes, max_memory_bytes))
        except Exception:
            pass

    return _apply_limits


def _excel_engine_for_extension(extension: str) -> Optional[str]:
    ext = _normalize_extension(extension)
    if ext == "xls":
        return "xlrd"
    if ext in {"xlsx", "xlsm"}:
        return "openpyxl"
    if ext == "xlsb":
        return "pyxlsb"
    return None


def _is_excel_placeholder_header(value: str) -> bool:
    normalized = normalize_text(value)
    if not normalized:
        return True

    lowered = normalized.lower()
    if lowered in {"nan", "none", "null"}:
        return True
    if lowered.startswith("unnamed:"):
        return True
    return False


def _make_unique_headers(headers: List[str]) -> List[str]:
    seen_counts: Dict[str, int] = {}
    unique_headers: List[str] = []

    for index, raw_header in enumerate(headers):
        header = normalize_text(raw_header) or f"Column {index + 1}"
        key = header.lower()
        seen_counts[key] = seen_counts.get(key, 0) + 1
        if seen_counts[key] == 1:
            unique_headers.append(header)
        else:
            unique_headers.append(f"{header}_{seen_counts[key]}")

    return unique_headers


def _clean_excel_headers_and_rows(dataframe: pd.DataFrame) -> tuple:
    if dataframe is None:
        return [], pd.DataFrame()

    cleaned_dataframe = dataframe.fillna("")
    raw_headers = [normalize_text("" if col is None else str(col)) for col in cleaned_dataframe.columns.tolist()]

    placeholder_count = sum(1 for header in raw_headers if _is_excel_placeholder_header(header))
    if raw_headers and placeholder_count >= max(1, len(raw_headers) // 2) and not cleaned_dataframe.empty:
        first_row_values = [normalize_text(str(value)) for value in cleaned_dataframe.iloc[0].tolist()]
        first_row_signal = sum(1 for value in first_row_values if value and not _is_excel_placeholder_header(value))
        if first_row_signal >= max(2, len(raw_headers) // 2):
            raw_headers = first_row_values
            cleaned_dataframe = cleaned_dataframe.iloc[1:].reset_index(drop=True)

    normalized_headers = []
    for index, header in enumerate(raw_headers):
        if _is_excel_placeholder_header(header):
            normalized_headers.append(f"Column {index + 1}")
        else:
            normalized_headers.append(header)

    normalized_headers = _make_unique_headers(normalized_headers)
    return normalized_headers, cleaned_dataframe


def _safe_int(value: Any) -> int:
    try:
        return int(value)
    except Exception:
        return 0


def _point_inside_bbox(x: float, y: float, bbox: tuple) -> bool:
    if not bbox or len(bbox) != 4:
        return False
    x0, top, x1, bottom = bbox
    return float(x0) <= x <= float(x1) and float(top) <= y <= float(bottom)


def _words_to_lines(words: List[Dict[str, Any]], y_tolerance: float = 3.0) -> str:
    if not words:
        return ""

    sorted_words = sorted(
        words,
        key=lambda word: (float(word.get("top", 0.0)), float(word.get("x0", 0.0))),
    )
    lines: List[str] = []
    current_words: List[str] = []
    current_top: Optional[float] = None

    for word in sorted_words:
        text = normalize_text(str(word.get("text", "")))
        if not text:
            continue

        word_top = float(word.get("top", 0.0))
        if current_top is None or abs(word_top - current_top) <= y_tolerance:
            current_words.append(text)
            if current_top is None:
                current_top = word_top
            else:
                current_top = (current_top + word_top) / 2.0
            continue

        if current_words:
            lines.append(" ".join(current_words))
        current_words = [text]
        current_top = word_top

    if current_words:
        lines.append(" ".join(current_words))

    return normalize_text("\n".join(lines))


def _extract_pdf_text_excluding_tables(page, table_bboxes: List[tuple]) -> str:
    try:
        words = page.extract_words(
            x_tolerance=2,
            y_tolerance=2,
            keep_blank_chars=False,
            use_text_flow=True,
        )
    except Exception:
        words = []

    if not words:
        return ""

    filtered_words = []
    for word in words:
        x0 = float(word.get("x0", 0.0))
        x1 = float(word.get("x1", x0))
        top = float(word.get("top", 0.0))
        bottom = float(word.get("bottom", top))
        mid_x = (x0 + x1) / 2.0
        mid_y = (top + bottom) / 2.0

        if any(_point_inside_bbox(mid_x, mid_y, bbox) for bbox in table_bboxes):
            continue
        filtered_words.append(word)

    return _words_to_lines(filtered_words)


def _table_signature(rows: List[List[Any]]) -> str:
    if not rows:
        return ""

    normalized_rows = []
    for row in rows[:40]:
        if not isinstance(row, (list, tuple)):
            continue
        normalized_cells = [normalize_text("" if cell is None else str(cell)) for cell in row[:24]]
        if any(normalized_cells):
            normalized_rows.append("|".join(normalized_cells))

    return "||".join(normalized_rows[:16])


def _extract_pdf_tables_with_fallback(page) -> tuple:
    settings_candidates = [
        {
            "vertical_strategy": "lines",
            "horizontal_strategy": "lines",
            "intersection_x_tolerance": 5,
            "intersection_y_tolerance": 5,
        },
        {
            "vertical_strategy": "lines",
            "horizontal_strategy": "text",
            "intersection_x_tolerance": 5,
            "intersection_y_tolerance": 5,
        },
        {
            "vertical_strategy": "text",
            "horizontal_strategy": "text",
            "snap_tolerance": 3,
            "join_tolerance": 3,
        },
    ]

    table_bboxes: List[tuple] = []
    table_rows_collection: List[List[List[str]]] = []
    seen_bboxes = set()
    seen_signatures = set()
    fallback_hits = 0

    for settings_index, table_settings in enumerate(settings_candidates):
        try:
            table_objects = page.find_tables(table_settings=table_settings)
        except Exception:
            table_objects = []

        if settings_index > 0 and table_objects:
            fallback_hits += 1

        for table in table_objects:
            bbox = getattr(table, "bbox", None)
            bbox_key = None
            if bbox and len(bbox) == 4:
                bbox_key = tuple(round(float(value), 1) for value in bbox)
                if bbox_key not in seen_bboxes:
                    table_bboxes.append(tuple(float(value) for value in bbox))
                    seen_bboxes.add(bbox_key)

            try:
                raw_rows = table.extract()
            except Exception:
                raw_rows = []

            normalized_rows = []
            for raw_row in raw_rows or []:
                if not isinstance(raw_row, (list, tuple)):
                    continue
                normalized_row = [normalize_text("" if cell is None else str(cell)) for cell in raw_row]
                normalized_rows.append(normalized_row)

            signature = _table_signature(normalized_rows)
            if signature and signature not in seen_signatures:
                seen_signatures.add(signature)
                table_rows_collection.append(normalized_rows)

    if table_rows_collection:
        return table_bboxes, table_rows_collection, fallback_hits

    for settings_index, table_settings in enumerate(settings_candidates):
        try:
            extracted_tables = page.extract_tables(table_settings=table_settings)
        except Exception:
            extracted_tables = []

        if settings_index > 0 and extracted_tables:
            fallback_hits += 1

        for raw_table in extracted_tables or []:
            normalized_rows = []
            for raw_row in raw_table or []:
                if not isinstance(raw_row, (list, tuple)):
                    continue
                normalized_row = [normalize_text("" if cell is None else str(cell)) for cell in raw_row]
                normalized_rows.append(normalized_row)

            signature = _table_signature(normalized_rows)
            if signature and signature not in seen_signatures:
                seen_signatures.add(signature)
                table_rows_collection.append(normalized_rows)

    return table_bboxes, table_rows_collection, fallback_hits


def _resize_image_for_ocr(image):
    width, height = image.size
    if width <= 0 or height <= 0:
        return image

    min_side = min(width, height)
    max_side = max(width, height)
    scale = 1.0

    if min_side < OCR_MIN_IMAGE_SIDE_PX:
        scale = OCR_MIN_IMAGE_SIDE_PX / float(min_side)
    if max_side * scale > OCR_MAX_IMAGE_SIDE_PX:
        scale = OCR_MAX_IMAGE_SIDE_PX / float(max_side)

    if abs(scale - 1.0) < 0.05:
        return image

    new_size = (
        max(1, int(round(width * scale))),
        max(1, int(round(height * scale))),
    )
    return image.resize(new_size)


def _prepare_image_for_ocr(raw_image):
    image = raw_image
    if ImageOps is not None:
        image = ImageOps.exif_transpose(image)

    if image.mode not in {"L", "RGB"}:
        image = image.convert("RGB")

    if OCR_IMAGE_PREPROCESS and ImageOps is not None:
        image = ImageOps.grayscale(image)
        image = ImageOps.autocontrast(image)

    image = _resize_image_for_ocr(image)
    return image


def _as_float(value: Any, default: float = 0.0) -> float:
    try:
        return float(value)
    except Exception:
        return default


def _polygon_to_bbox(box: Any) -> Optional[tuple]:
    if not isinstance(box, (list, tuple)) or not box:
        return None

    points = []
    for point in box:
        if not isinstance(point, (list, tuple)) or len(point) < 2:
            continue
        x = _as_float(point[0], default=float("nan"))
        y = _as_float(point[1], default=float("nan"))
        if x != x or y != y:
            continue
        points.append((x, y))

    if not points:
        return None

    xs = [point[0] for point in points]
    ys = [point[1] for point in points]
    return min(xs), min(ys), max(xs), max(ys)


def _polygon_anchor(box: Any) -> tuple:
    bbox = _polygon_to_bbox(box)
    if bbox is not None:
        x0, top, _, _ = bbox
        return float(top), float(x0)

    if isinstance(box, (list, tuple)) and box:
        first_point = box[0]
        if isinstance(first_point, (list, tuple)) and len(first_point) >= 2:
            return float(first_point[1]), float(first_point[0])
    return 0.0, 0.0


def _canonical_text_for_similarity(text: str) -> str:
    normalized = normalize_text(text).lower()
    if not normalized:
        return ""

    folded = unicodedata.normalize("NFKD", normalized)
    folded = "".join(ch for ch in folded if not unicodedata.combining(ch))
    folded = re.sub(r"[^a-z0-9]+", " ", folded)
    return re.sub(r"\s+", " ", folded).strip()


def _similarity_ratio(left: str, right: str) -> float:
    if not left or not right:
        return 0.0
    if left == right:
        return 1.0

    short, long = (left, right) if len(left) <= len(right) else (right, left)
    if short and short in long:
        if len(short) >= 24:
            return len(short) / max(1, len(long))

    return SequenceMatcher(None, left, right).ratio()


def _line_exists_with_similarity(
    candidate_line: str,
    canonical_existing_lines: List[str],
    similarity_threshold: float,
) -> bool:
    canonical_candidate = _canonical_text_for_similarity(candidate_line)
    if not canonical_candidate:
        return True

    if canonical_candidate in canonical_existing_lines:
        return True

    for existing in canonical_existing_lines:
        length_gap = abs(len(existing) - len(canonical_candidate))
        max_len = max(len(existing), len(canonical_candidate), 1)
        if length_gap > max(8, int(round(max_len * 0.45))):
            continue
        if _similarity_ratio(existing, canonical_candidate) >= similarity_threshold:
            return True

    return False


def _cluster_numeric_values(values: List[float], tolerance: float) -> List[float]:
    if not values:
        return []

    sorted_values = sorted(values)
    clusters: List[List[float]] = [[sorted_values[0]]]

    for value in sorted_values[1:]:
        if abs(value - clusters[-1][-1]) <= tolerance:
            clusters[-1].append(value)
        else:
            clusters.append([value])

    anchors = [sum(cluster) / len(cluster) for cluster in clusters if cluster]
    return sorted(anchors)


def _cluster_blocks_by_rows(ocr_blocks: List[Dict[str, Any]], y_tolerance: float) -> List[List[Dict[str, Any]]]:
    if not ocr_blocks:
        return []

    ordered = sorted(ocr_blocks, key=lambda block: (float(block.get("top", 0.0)), float(block.get("x0", 0.0))))
    rows: List[Dict[str, Any]] = []

    for block in ordered:
        block_top = _as_float(block.get("top"), 0.0)
        block_bottom = _as_float(block.get("bottom"), block_top)
        assigned_row: Optional[Dict[str, Any]] = None

        for row in rows:
            overlaps_vertically = (
                block_top <= float(row["bottom"]) + y_tolerance
                and block_bottom >= float(row["top"]) - y_tolerance
            )
            near_center = abs(block_top - float(row["center"])) <= y_tolerance
            if overlaps_vertically or near_center:
                assigned_row = row
                break

        if assigned_row is None:
            rows.append(
                {
                    "items": [block],
                    "top": block_top,
                    "bottom": block_bottom,
                    "center": block_top,
                }
            )
            continue

        assigned_row["items"].append(block)
        assigned_row["top"] = min(float(assigned_row["top"]), block_top)
        assigned_row["bottom"] = max(float(assigned_row["bottom"]), block_bottom)
        item_count = len(assigned_row["items"])
        assigned_row["center"] = ((float(assigned_row["center"]) * (item_count - 1)) + block_top) / float(
            item_count
        )

    finalized_rows = []
    for row in rows:
        row_items = sorted(row["items"], key=lambda item: float(item.get("x0", 0.0)))
        if row_items:
            finalized_rows.append(row_items)

    finalized_rows.sort(key=lambda row_items: float(row_items[0].get("top", 0.0)))
    return finalized_rows


def _parse_rapidocr_result(result: Any) -> tuple:
    lines_with_anchor = []
    scores = []
    ocr_blocks: List[Dict[str, Any]] = []

    if not isinstance(result, list):
        return "", 0.0, 0, ocr_blocks

    for item in result:
        if not isinstance(item, (list, tuple)) or len(item) < 3:
            continue

        text = normalize_text(str(item[1]))
        if not text:
            continue

        try:
            score = float(item[2])
        except Exception:
            score = 0.0

        if score < OCR_MIN_CONFIDENCE:
            continue

        bbox = _polygon_to_bbox(item[0])
        anchor_y, anchor_x = _polygon_anchor(item[0])
        if bbox is None:
            x0 = anchor_x
            top = anchor_y
            x1 = anchor_x
            bottom = anchor_y
        else:
            x0, top, x1, bottom = bbox

        lines_with_anchor.append((anchor_y, anchor_x, text))
        scores.append(score)
        ocr_blocks.append(
            {
                "text": text,
                "score": score,
                "x0": float(x0),
                "top": float(top),
                "x1": float(x1),
                "bottom": float(bottom),
            }
        )

    if not lines_with_anchor:
        return "", 0.0, 0, []

    lines_with_anchor.sort(key=lambda item: (item[0], item[1]))
    ordered_lines = [line for _, _, line in lines_with_anchor]
    avg_score = sum(scores) / len(scores) if scores else 0.0
    ordered_blocks = sorted(ocr_blocks, key=lambda block: (float(block["top"]), float(block["x0"])))
    return normalize_text("\n".join(ordered_lines)), avg_score, len(ordered_lines), ordered_blocks


def _empty_ocr_layout_metadata() -> Dict[str, Any]:
    return {
        "box_count": 0,
        "table_added": False,
        "diagram_added": False,
        "table_rows": 0,
        "table_cols": 0,
        "table_density": 0.0,
    }


def _reconstruct_markdown_table_from_ocr_blocks(ocr_blocks: List[Dict[str, Any]]) -> tuple:
    if not OCR_TABLE_RECONSTRUCTION_ENABLED:
        return "", _empty_ocr_layout_metadata()

    text_blocks = [block for block in ocr_blocks if normalize_text(str(block.get("text", "")))]
    if len(text_blocks) < OCR_TABLE_MIN_BOXES:
        return "", _empty_ocr_layout_metadata()

    widths = [
        max(1.0, float(block.get("x1", 0.0)) - float(block.get("x0", 0.0)))
        for block in text_blocks
        if float(block.get("x1", 0.0)) >= float(block.get("x0", 0.0))
    ]
    heights = [
        max(1.0, float(block.get("bottom", 0.0)) - float(block.get("top", 0.0)))
        for block in text_blocks
        if float(block.get("bottom", 0.0)) >= float(block.get("top", 0.0))
    ]
    if not widths or not heights:
        return "", _empty_ocr_layout_metadata()

    y_tolerance = max(5.0, statistics.median(heights) * 0.75)
    rows = _cluster_blocks_by_rows(text_blocks, y_tolerance=y_tolerance)
    if len(rows) < 3:
        return "", _empty_ocr_layout_metadata()

    x_centers = [
        (float(block.get("x0", 0.0)) + float(block.get("x1", 0.0))) / 2.0 for block in text_blocks
    ]
    x_tolerance = max(8.0, statistics.median(widths) * 0.65)
    column_anchors = _cluster_numeric_values(x_centers, tolerance=x_tolerance)
    if len(column_anchors) < OCR_TABLE_MIN_COLUMNS or len(column_anchors) > 14:
        return "", _empty_ocr_layout_metadata()

    row_values: List[List[str]] = []
    for row in rows[:120]:
        cells = ["" for _ in column_anchors]
        for block in sorted(row, key=lambda item: float(item.get("x0", 0.0))):
            center_x = (float(block.get("x0", 0.0)) + float(block.get("x1", 0.0))) / 2.0
            col_idx = min(range(len(column_anchors)), key=lambda idx: abs(column_anchors[idx] - center_x))
            text = normalize_text(str(block.get("text", "")))
            if not text:
                continue

            if cells[col_idx]:
                if text not in cells[col_idx]:
                    cells[col_idx] = f"{cells[col_idx]} {text}"
            else:
                cells[col_idx] = text

        if any(cells):
            row_values.append(cells)

    if len(row_values) < 3:
        return "", _empty_ocr_layout_metadata()

    filled_per_row = [sum(1 for cell in row if cell) for row in row_values]
    if max(filled_per_row, default=0) < OCR_TABLE_MIN_COLUMNS:
        return "", _empty_ocr_layout_metadata()

    non_empty_cells = sum(1 for row in row_values for cell in row if cell)
    density = non_empty_cells / float(max(1, len(row_values) * len(column_anchors)))
    if density < OCR_TABLE_MIN_DENSITY:
        return "", _empty_ocr_layout_metadata()

    header_index = 0
    for idx, row in enumerate(row_values[:4]):
        if sum(1 for value in row if value) >= max(2, len(row) // 2):
            header_index = idx
            break

    if header_index > 0:
        row_values = row_values[header_index:]

    markdown_table = rows_to_markdown_table(row_values)
    if not markdown_table:
        return "", _empty_ocr_layout_metadata()

    metadata = _empty_ocr_layout_metadata()
    metadata["box_count"] = len(text_blocks)
    metadata["table_added"] = True
    metadata["table_rows"] = len(row_values)
    metadata["table_cols"] = len(column_anchors)
    metadata["table_density"] = float(round(density, 3))
    return markdown_table, metadata


def _build_diagram_summary_from_ocr_blocks(ocr_blocks: List[Dict[str, Any]]) -> tuple:
    if not OCR_DIAGRAM_SUMMARY_ENABLED:
        return "", _empty_ocr_layout_metadata()

    text_blocks = [block for block in ocr_blocks if normalize_text(str(block.get("text", "")))]
    if len(text_blocks) < OCR_DIAGRAM_MIN_BOXES:
        return "", _empty_ocr_layout_metadata()

    heights = [
        max(1.0, float(block.get("bottom", 0.0)) - float(block.get("top", 0.0)))
        for block in text_blocks
        if float(block.get("bottom", 0.0)) >= float(block.get("top", 0.0))
    ]
    if not heights:
        return "", _empty_ocr_layout_metadata()

    y_tolerance = max(6.0, statistics.median(heights) * 0.8)
    rows = _cluster_blocks_by_rows(text_blocks, y_tolerance=y_tolerance)
    if len(rows) < 3:
        return "", _empty_ocr_layout_metadata()

    avg_label_length = sum(len(str(block.get("text", ""))) for block in text_blocks) / float(len(text_blocks))
    if avg_label_length > 24:
        return "", _empty_ocr_layout_metadata()

    summary_lines = ["[DIAGRAM_START]", "Diagram OCR labels (layout-aware reconstruction):"]
    labels_count = 0
    for row_idx, row in enumerate(rows[:10], start=1):
        labels = []
        for block in row:
            label = normalize_text(str(block.get("text", "")))
            if label:
                labels.append(label)
        if not labels:
            continue
        labels_count += len(labels)
        compact_labels = " | ".join(labels[:8])
        summary_lines.append(f"- lane_{row_idx}: {compact_labels}")

    if labels_count < OCR_DIAGRAM_MIN_BOXES:
        return "", _empty_ocr_layout_metadata()

    summary_lines.append("[DIAGRAM_END]")
    metadata = _empty_ocr_layout_metadata()
    metadata["box_count"] = len(text_blocks)
    metadata["diagram_added"] = True
    return "\n".join(summary_lines), metadata


def _augment_ocr_text_with_layout(ocr_text: str, ocr_blocks: List[Dict[str, Any]]) -> tuple:
    base_text = normalize_text(ocr_text)
    if not ocr_blocks:
        return base_text, _empty_ocr_layout_metadata()

    merged_sections = [base_text] if base_text else []
    table_text, table_metadata = _reconstruct_markdown_table_from_ocr_blocks(ocr_blocks)
    merged_metadata = _empty_ocr_layout_metadata()
    merged_metadata["box_count"] = len(ocr_blocks)

    if table_text:
        merged_sections.append(f"[TABLE_START]\n{table_text}\n[TABLE_END]")
        merged_metadata.update(table_metadata)
        merged_metadata["box_count"] = len(ocr_blocks)
    else:
        diagram_summary, diagram_metadata = _build_diagram_summary_from_ocr_blocks(ocr_blocks)
        if diagram_summary:
            merged_sections.append(diagram_summary)
            merged_metadata.update(diagram_metadata)
            merged_metadata["box_count"] = len(ocr_blocks)

    return normalize_text("\n\n".join(section for section in merged_sections if section)), merged_metadata


def _get_rapidocr_engine() -> Optional[Any]:
    global _RAPID_OCR_ENGINE

    if RapidOCR is None:
        return None
    if _RAPID_OCR_ENGINE is not None:
        return _RAPID_OCR_ENGINE

    try:
        _RAPID_OCR_ENGINE = RapidOCR()
    except Exception:
        _RAPID_OCR_ENGINE = None
    return _RAPID_OCR_ENGINE


def _resolve_tesseract_binary() -> str:
    if os.path.isabs(OCR_TESSERACT_BIN) and os.path.exists(OCR_TESSERACT_BIN):
        return OCR_TESSERACT_BIN
    return shutil.which(OCR_TESSERACT_BIN) or OCR_TESSERACT_BIN


def _ocr_with_rapidocr(image) -> tuple:
    if np is None:
        return "", "rapidocr:missing_numpy", []

    engine = _get_rapidocr_engine()
    if engine is None:
        return "", "rapidocr:unavailable", []

    try:
        result, _ = engine(np.asarray(image))
        text, avg_score, line_count, ocr_blocks = _parse_rapidocr_result(result)
        if not text:
            return "", "rapidocr:no_text", ocr_blocks
        return text, f"rapidocr|lines:{line_count}|avg_conf:{avg_score:.2f}", ocr_blocks
    except Exception as exc:
        return "", f"rapidocr:error:{exc.__class__.__name__}", []


def _ocr_with_tesseract(image) -> tuple:
    if pytesseract is None:
        return "", "tesseract:module_missing", []

    binary = _resolve_tesseract_binary()
    if not binary or (not os.path.exists(binary) and not shutil.which(binary)):
        return "", "tesseract:binary_missing", []

    try:
        pytesseract.pytesseract.tesseract_cmd = binary
        config = f"--oem {OCR_TESSERACT_OEM} --psm {OCR_TESSERACT_PSM}"
        text = pytesseract.image_to_string(
            image,
            lang=OCR_TESSERACT_LANG,
            config=config,
            timeout=OCR_TIMEOUT_SECONDS,
        )
        normalized = normalize_text(text)
        if not normalized:
            return "", "tesseract:no_text", []
        return normalized, "tesseract", []
    except Exception as exc:
        return "", f"tesseract:error:{exc.__class__.__name__}", []


def _ocr_image_with_priority(image) -> Dict[str, Any]:
    traces: List[str] = []
    for engine in OCR_ENGINE_PRIORITY:
        if engine == "rapidocr":
            text, trace, ocr_blocks = _ocr_with_rapidocr(image)
        elif engine == "tesseract":
            text, trace, ocr_blocks = _ocr_with_tesseract(image)
        else:
            text, trace, ocr_blocks = "", f"{engine}:unsupported", []

        traces.append(trace)
        if text:
            return {
                "text": text,
                "trace": ";".join(traces),
                "engine": engine,
                "blocks": ocr_blocks,
            }

    return {
        "text": "",
        "trace": ";".join(traces) if traces else "ocr:no_engine",
        "engine": None,
        "blocks": [],
    }


def _merge_native_and_ocr_text(native_text: str, ocr_text: str) -> str:
    native = normalize_text(native_text)
    ocr = normalize_text(ocr_text)

    if not native:
        return ocr
    if not ocr:
        return native

    canonical_native = _canonical_text_for_similarity(native)
    canonical_ocr = _canonical_text_for_similarity(ocr)
    if canonical_native and canonical_ocr:
        full_similarity = _similarity_ratio(canonical_native, canonical_ocr)
        if full_similarity >= OCR_SUPPLEMENT_TEXT_SIMILARITY:
            return native if len(native) >= len(ocr) else ocr

    if ocr in native:
        return native
    if native in ocr:
        return ocr

    native_lines = [normalize_text(line) for line in native.splitlines() if normalize_text(line)]
    canonical_native_lines = []
    for line in native_lines:
        canonical_line = _canonical_text_for_similarity(line)
        if canonical_line:
            canonical_native_lines.append(canonical_line)

    supplemental_lines: List[str] = []
    seen_supplement_canonical: List[str] = []
    for line in [normalize_text(item) for item in ocr.splitlines() if normalize_text(item)]:
        canonical_line = _canonical_text_for_similarity(line)
        if not canonical_line:
            continue
        if canonical_line in seen_supplement_canonical:
            continue
        if _line_exists_with_similarity(line, canonical_native_lines, OCR_SUPPLEMENT_LINE_SIMILARITY):
            continue

        supplemental_lines.append(line)
        seen_supplement_canonical.append(canonical_line)
        canonical_native_lines.append(canonical_line)

    if not supplemental_lines:
        return native

    return f"{native}\n[OCR_SUPPLEMENT]\n" + "\n".join(supplemental_lines)


def _extract_ocr_text_from_pdf_page(pdfium_document: Any, page_index: int) -> Dict[str, Any]:
    if pdfium_document is None:
        return {
            "text": "",
            "trace": "ocr_pdf:pdfium_unavailable",
            "engine": None,
            "blocks": [],
            "layout": _empty_ocr_layout_metadata(),
        }
    if Image is None:
        return {
            "text": "",
            "trace": "ocr_pdf:pillow_unavailable",
            "engine": None,
            "blocks": [],
            "layout": _empty_ocr_layout_metadata(),
        }

    page = None
    bitmap = None
    try:
        page = pdfium_document[page_index]
        bitmap = page.render(scale=OCR_PDF_RENDER_DPI / 72.0)
        image = bitmap.to_pil()
        prepared = _prepare_image_for_ocr(image)
        ocr_result = _ocr_image_with_priority(prepared)
        enhanced_text, layout_metadata = _augment_ocr_text_with_layout(
            str(ocr_result.get("text") or ""),
            list(ocr_result.get("blocks") or []),
        )
        ocr_result["text"] = enhanced_text
        ocr_result["layout"] = layout_metadata
        return ocr_result
    except Exception as exc:
        return {
            "text": "",
            "trace": f"ocr_pdf:error:{exc.__class__.__name__}",
            "engine": None,
            "blocks": [],
            "layout": _empty_ocr_layout_metadata(),
        }
    finally:
        try:
            if page is not None:
                page.close()
        except Exception:
            pass
        try:
            if bitmap is not None:
                bitmap.close()
        except Exception:
            pass


def _flatten_json(prefix: str, node: Any, out: List[str]) -> None:
    if isinstance(node, dict):
        for key, value in node.items():
            child_prefix = f"{prefix}.{key}" if prefix else str(key)
            _flatten_json(child_prefix, value, out)
        return

    if isinstance(node, list):
        for index, value in enumerate(node):
            child_prefix = f"{prefix}[{index}]"
            _flatten_json(child_prefix, value, out)
        return

    value = "" if node is None else str(node)
    normalized = normalize_text(value)
    if normalized:
        out.append(f"{prefix}: {normalized}" if prefix else normalized)


def _flatten_xml(root: etree._Element, out: List[str]) -> None:
    for element in root.iter():
        tag_name = str(element.tag)
        path = root.getroottree().getpath(element)
        if element.attrib:
            for key, value in element.attrib.items():
                normalized_value = normalize_text(str(value))
                if normalized_value:
                    out.append(f"{path}.@{key}: {normalized_value}")
        text_value = normalize_text(element.text or "")
        if text_value:
            out.append(f"{path}<{tag_name}>: {text_value}")


def _extract_html_tables_markdown(html_text: str, max_tables: int = 8, max_rows: int = 40) -> List[str]:
    markdown_tables: List[str] = []
    if not html_text:
        return markdown_tables

    try:
        dataframes = pd.read_html(io.StringIO(html_text))
    except Exception:
        return markdown_tables

    for dataframe in dataframes[:max_tables]:
        try:
            cleaned_dataframe = dataframe.fillna("")
            raw_headers = [normalize_text("" if col is None else str(col)) for col in cleaned_dataframe.columns.tolist()]

            normalized_headers = []
            for header_index, header_value in enumerate(raw_headers):
                if _is_excel_placeholder_header(header_value):
                    normalized_headers.append(f"Column {header_index + 1}")
                else:
                    normalized_headers.append(header_value)
            normalized_headers = _make_unique_headers(normalized_headers)

            preview_rows = [normalized_headers]
            for row_index, row_values in enumerate(cleaned_dataframe.itertuples(index=False, name=None)):
                if row_index >= max_rows:
                    break
                preview_rows.append([normalize_text("" if value is None else str(value)) for value in row_values])

            markdown_preview = rows_to_markdown_table(preview_rows)
            if markdown_preview:
                markdown_tables.append(f"[TABLE_START]\n{markdown_preview}\n[TABLE_END]")
        except Exception:
            continue

    return markdown_tables


def extract_text_from_textlike(content: bytes, extension: str) -> Dict[str, Any]:
    ext = _normalize_extension(extension)
    strategy = _with_lib_versions(f"textlike:{ext}", "beautifulsoup4", "lxml", "pandas")
    if not content:
        return _response("", _error("ERR_EMPTY", "Empty"), strategy)

    decoded = safe_decode(content)
    if ext in {"html", "htm"}:
        raw_html = decoded
        decoded = strip_html_tags(raw_html)
        html_tables = _extract_html_tables_markdown(raw_html)
        if html_tables:
            decoded = normalize_text("\n\n".join([decoded, *html_tables]))
            strategy = f"{strategy}|html_tables:{len(html_tables)}"
    elif ext == "rtf":
        decoded = _strip_rtf_control_words(decoded)
        strategy = f"{strategy}|rtf_stripped"
    elif ext == "json":
        try:
            parsed_json = json.loads(decoded)
            flattened_lines: List[str] = []
            _flatten_json("", parsed_json, flattened_lines)
            decoded = "\n".join(flattened_lines)
            strategy = f"{strategy}|flattened"
        except Exception as exc:
            strategy = f"{strategy}|json_fallback:{exc.__class__.__name__}"
    elif ext == "xml":
        try:
            root = etree.fromstring(decoded.encode("utf-8", errors="ignore"))
            flattened_lines: List[str] = []
            _flatten_xml(root, flattened_lines)
            decoded = "\n".join(flattened_lines)
            strategy = f"{strategy}|flattened"
        except Exception as exc:
            strategy = f"{strategy}|xml_fallback:{exc.__class__.__name__}"
    elif ext in {"csv", "tsv"}:
        try:
            dataframe = pd.read_csv(
                io.StringIO(decoded),
                sep="\t" if ext == "tsv" else ",",
                dtype=str,
                engine="python",
                quotechar='"',
                escapechar="\\",
                on_bad_lines="skip",
                keep_default_na=False,
                na_filter=False,
                nrows=MAX_CSV_ROWS,
            )

            cleaned_dataframe = dataframe.fillna("")
            raw_headers = [normalize_text("" if col is None else str(col)) for col in cleaned_dataframe.columns.tolist()]
            normalized_headers = []
            for header_index, header_value in enumerate(raw_headers):
                if _is_excel_placeholder_header(header_value):
                    normalized_headers.append(f"Column {header_index + 1}")
                else:
                    normalized_headers.append(header_value)
            normalized_headers = _make_unique_headers(normalized_headers)

            lines: List[str] = []
            preview_rows = [normalized_headers]

            for row_index, row_values in enumerate(cleaned_dataframe.itertuples(index=False, name=None)):
                normalized_values = [normalize_text("" if value is None else str(value)) for value in row_values]

                if row_index < MAX_CSV_TABLE_PREVIEW_ROWS:
                    preview_rows.append(normalized_values)

                if EXCEL_INCLUDE_ROW_CONTEXT:
                    row_context = ", ".join(
                        [
                            f"{str(header)}: {str(value)}"
                            for header, value in zip(normalized_headers, normalized_values)
                            if value
                        ]
                    )
                    if row_context:
                        lines.append(f"- {row_context}")

            markdown_preview = rows_to_markdown_table(preview_rows)
            if markdown_preview:
                lines.insert(0, f"[TABLE_START]\n{markdown_preview}\n[TABLE_END]")

            if len(cleaned_dataframe) >= MAX_CSV_ROWS:
                lines.append(f"[WARNING] CSV row limit reached ({MAX_CSV_ROWS} rows)")

            decoded = "\n".join(lines) if lines else decoded
            strategy = f"{strategy}|table_preview_rows:{MAX_CSV_TABLE_PREVIEW_ROWS}"
        except Exception as exc:
            decoded_preview = _truncate_text(decoded)
            decoded = f"### CSV Data\n[TABLE_START]\n{decoded_preview}\n[TABLE_END]"
            strategy = f"{strategy}|csv_fallback:{exc.__class__.__name__}"

    decoded = normalize_text(decoded)
    return _response(decoded, None if decoded else _error("ERR_EMPTY", "Empty decoded"), strategy)


def extract_text_from_image_ocr(content: bytes, extension: str) -> Dict[str, Any]:
    ext = _normalize_extension(extension)
    strategy = _with_lib_versions(
        f"ocr_image:{ext}|engines:{','.join(OCR_ENGINE_PRIORITY)}",
        "pillow",
        "rapidocr-onnxruntime",
        "pytesseract",
    )

    if not content:
        return _response("", _error("ERR_EMPTY", "Empty"), strategy)
    if not OCR_ENABLED:
        return _response("", _error("ERR_OCR_DISABLED", "OCR disabled by configuration"), f"{strategy}|disabled")
    if Image is None:
        return _response("", _error("ERR_OCR_DEPENDENCY", "Pillow unavailable"), strategy)

    try:
        with Image.open(io.BytesIO(content)) as image_obj:
            prepared = _prepare_image_for_ocr(image_obj.copy())
    except Exception as exc:
        return _response("", _error("ERR_OCR_IMAGE", str(exc)), strategy)

    ocr_result = _ocr_image_with_priority(prepared)
    enhanced_text, layout_metadata = _augment_ocr_text_with_layout(
        str(ocr_result.get("text") or ""),
        list(ocr_result.get("blocks") or []),
    )
    ocr_trace_base = str(ocr_result.get("trace") or "ocr:no_engine")
    ocr_trace = (
        f"{ocr_trace_base}|layout_boxes:{int(layout_metadata.get('box_count', 0))}"
        f"|layout_table:{int(bool(layout_metadata.get('table_added')))}"
        f"|layout_diagram:{int(bool(layout_metadata.get('diagram_added')))}"
    )

    if enhanced_text:
        return _with_ocr_metadata(
            _response(enhanced_text, None, f"{strategy}|{ocr_trace}"),
            attempted=True,
            used=True,
            trace=ocr_trace,
            pages=1,
            supplement_pages=1,
        )

    return _with_ocr_metadata(
        _response("", _error("ERR_OCR_EMPTY", "No text recognized"), f"{strategy}|{ocr_trace}"),
        attempted=True,
        used=False,
        trace=ocr_trace,
        pages=0,
        supplement_pages=0,
    )


def extract_text_from_pdf_pdfplumber(content: bytes, extension: str) -> Dict[str, Any]:
    ext = _normalize_extension(extension) or "pdf"
    strategy = _with_lib_versions(
        f"pdfplumber:{ext}|table_text_excluded:{int(PDF_EXCLUDE_TABLE_TEXT)}",
        "pdfplumber",
    )
    if not content:
        return _response("", _error("ERR_EMPTY", "Empty"), strategy)

    try:
        started_at = time.time()
        elements = []
        pdfium_document = None
        ocr_pages = 0
        ocr_supplement_pages = 0
        ocr_attempted_pages = 0
        ocr_traces: List[str] = []
        ocr_layout_table_pages = 0
        ocr_layout_diagram_pages = 0
        table_fallback_hits = 0
        ocr_attempt_budget_hit = False
        ocr_time_budget_hit = False

        if OCR_ENABLED and OCR_PDF_FALLBACK_ENABLED and pdfium is not None:
            try:
                pdfium_document = pdfium.PdfDocument(content)
            except Exception as exc:
                strategy = f"{strategy}|ocr_pdfium_init:{exc.__class__.__name__}"

        with pdfplumber.open(io.BytesIO(content)) as pdf:
            total_pages = len(pdf.pages)
            page_limit_reached = total_pages > MAX_PDF_PAGES
            pages = pdf.pages[:MAX_PDF_PAGES]
            used_layout_fallback = False
            timeout_partial = False
            last_page_processed = 0

            if page_limit_reached:
                elements.append(f"[WARNING] PDF page limit reached ({MAX_PDF_PAGES}/{total_pages})")

            for page in pages:
                if _is_timed_out(started_at):
                    timeout_partial = True
                    elements.append(
                        f"[WARNING] PDF parsing timeout reached after {PARSE_TIMEOUT_SECONDS}s at page {last_page_processed}. Returning partial extraction."
                    )
                    break

                elements.append(f"## Page {page.page_number}")
                last_page_processed = page.page_number

                working_page = page
                try:
                    working_page = page.dedupe_chars(tolerance=1)
                except Exception:
                    pass

                table_bboxes, extracted_tables, fallback_hits = _extract_pdf_tables_with_fallback(working_page)
                if fallback_hits > 0:
                    table_fallback_hits += 1

                page_text = ""
                if PDF_EXCLUDE_TABLE_TEXT and table_bboxes:
                    page_text = _extract_pdf_text_excluding_tables(working_page, table_bboxes)

                if not page_text:
                    page_text = working_page.extract_text(layout=False, x_tolerance=2, y_tolerance=2)
                if not page_text:
                    page_text = working_page.extract_text(layout=True, x_tolerance=2, y_tolerance=2)
                    if page_text:
                        used_layout_fallback = True

                native_char_count = len(normalize_text(page_text or ""))
                ocr_elapsed_seconds = time.time() - started_at
                ocr_time_budget_reached = ocr_elapsed_seconds >= (PARSE_TIMEOUT_SECONDS * OCR_PDF_OCR_TIME_BUDGET_RATIO)
                ocr_attempt_budget_reached = ocr_attempted_pages >= OCR_MAX_PDF_OCR_ATTEMPTS

                ocr_eligible = (
                    OCR_ENABLED
                    and OCR_PDF_FALLBACK_ENABLED
                    and page.page_number <= OCR_MAX_PDF_PAGES
                    and native_char_count < OCR_MIN_PDF_PAGE_TEXT_CHARS
                )

                if ocr_eligible and ocr_time_budget_reached:
                    ocr_time_budget_hit = True
                if ocr_eligible and ocr_attempt_budget_reached:
                    ocr_attempt_budget_hit = True

                should_run_ocr = (
                    ocr_eligible
                    and not ocr_time_budget_reached
                    and not ocr_attempt_budget_reached
                )

                if should_run_ocr:
                    ocr_attempted_pages += 1
                    ocr_result = _extract_ocr_text_from_pdf_page(pdfium_document, page.page_number - 1)
                    ocr_text = normalize_text(str(ocr_result.get("text") or ""))
                    ocr_trace = str(ocr_result.get("trace") or "ocr:no_trace")
                    layout_metadata = dict(ocr_result.get("layout") or {})

                    if layout_metadata.get("table_added"):
                        ocr_layout_table_pages += 1
                    if layout_metadata.get("diagram_added"):
                        ocr_layout_diagram_pages += 1

                    ocr_traces.append(f"p{page.page_number}:{ocr_trace}")
                    if ocr_text:
                        ocr_pages += 1
                        if page_text:
                            merged_page_text = _merge_native_and_ocr_text(page_text, ocr_text)
                            if merged_page_text != page_text:
                                ocr_supplement_pages += 1
                            page_text = merged_page_text
                        else:
                            page_text = ocr_text
                            ocr_supplement_pages += 1

                if page_text:
                    elements.append(page_text)

                for table in extracted_tables:
                    markdown_table = rows_to_markdown_table(table)
                    if markdown_table:
                        elements.append(f"\n[TABLE_START]\n{markdown_table}\n[TABLE_END]\n")

            if used_layout_fallback:
                strategy = f"{strategy}|layout_fallback"
            if page_limit_reached:
                strategy = f"{strategy}|page_limit:{MAX_PDF_PAGES}"
            if timeout_partial:
                strategy = f"{strategy}|timeout_partial:{last_page_processed}"
            if OCR_ENABLED and OCR_PDF_FALLBACK_ENABLED:
                strategy = (
                    f"{strategy}|ocr_pages:{ocr_pages}|ocr_supplement_pages:{ocr_supplement_pages}|"
                    f"ocr_attempted_pages:{ocr_attempted_pages}"
                )
                strategy = (
                    f"{strategy}|ocr_layout_tables:{ocr_layout_table_pages}|"
                    f"ocr_layout_diagrams:{ocr_layout_diagram_pages}"
                )
                if ocr_attempt_budget_hit:
                    strategy = f"{strategy}|ocr_attempt_budget:{OCR_MAX_PDF_OCR_ATTEMPTS}"
                if ocr_time_budget_hit:
                    strategy = f"{strategy}|ocr_time_budget_ratio:{OCR_PDF_OCR_TIME_BUDGET_RATIO:.2f}"
                if ocr_traces:
                    strategy = f"{strategy}|ocr_trace:{','.join(ocr_traces[:6])}"
            if table_fallback_hits > 0:
                strategy = f"{strategy}|table_detector_fallback_pages:{table_fallback_hits}"

        try:
            if pdfium_document is not None:
                pdfium_document.close()
        except Exception:
            pass

        full_text = normalize_text("\n\n".join(elements))
        ocr_trace_joined = ",".join(ocr_traces[:20]) if ocr_traces else None
        if not full_text and timeout_partial:
            return _with_ocr_metadata(
                _response("", _error("ERR_TIMEOUT", f"PDF parsing exceeded {PARSE_TIMEOUT_SECONDS}s"), strategy),
                attempted=ocr_attempted_pages > 0,
                used=ocr_pages > 0,
                trace=ocr_trace_joined,
                pages=ocr_pages,
                supplement_pages=ocr_supplement_pages,
            )
        return _with_ocr_metadata(
            _response(full_text, None if full_text else _error("ERR_EMPTY", "Empty PDF"), strategy),
            attempted=ocr_attempted_pages > 0,
            used=ocr_pages > 0,
            trace=ocr_trace_joined,
            pages=ocr_pages,
            supplement_pages=ocr_supplement_pages,
        )
    except Exception as exc:
        lowered_message = str(exc).lower()
        if "password" in lowered_message or "encrypted" in lowered_message:
            return _response("", _error("ERR_PDF_PASSWORD", str(exc)), strategy)
        return _response("", _error("ERR_PDF", str(exc)), strategy)


def extract_text_from_docx(content: bytes, extension: str) -> Dict[str, Any]:
    ext = _normalize_extension(extension) or "docx"
    strategy = _with_lib_versions(f"python_docx:{ext}", "python-docx")
    if not content:
        return _response("", _error("ERR_EMPTY", "Empty"), strategy)

    try:
        started_at = time.time()
        document = Document(io.BytesIO(content))
        elements = []
        ocr_attempted_units = 0
        ocr_used_units = 0
        ocr_traces: List[str] = []
        ocr_layout_tables = 0
        ocr_layout_diagrams = 0

        for section_index, section in enumerate(document.sections, start=1):
            header_text = normalize_text("\n".join([p.text for p in section.header.paragraphs if p.text.strip()]))
            footer_text = normalize_text("\n".join([p.text for p in section.footer.paragraphs if p.text.strip()]))
            if header_text:
                elements.append(f"## Header {section_index}\n{header_text}")
            if footer_text:
                elements.append(f"## Footer {section_index}\n{footer_text}")

        comments = _extract_docx_comments(document)
        if comments:
            elements.append("## Comments")
            elements.extend([f"- {comment}" for comment in comments])

        for child in document.element.body.iterchildren():
            if _is_timed_out(started_at):
                return _response("", _error("ERR_TIMEOUT", f"DOCX parsing exceeded {PARSE_TIMEOUT_SECONDS}s"), strategy)

            if child.tag.endswith("p"):
                from docx.text.paragraph import Paragraph

                paragraph = Paragraph(child, document)
                style_name = paragraph.style.name.lower() if paragraph.style else ""
                text = _extract_paragraph_with_hyperlinks(paragraph)
                if text:
                    if style_name.startswith("heading"):
                        level = 1
                        parts = style_name.split()
                        if len(parts) > 1 and parts[1].isdigit():
                            level = max(1, min(6, int(parts[1])))
                        text = f"{'#' * level} {text}"
                    elif "list" in style_name:
                        text = f"- {text}"
                    elements.append(text)
            elif child.tag.endswith("tbl"):
                from docx.table import Table

                rows = [[cell.text.strip() for cell in row.cells] for row in Table(child, document).rows]
                markdown_table = rows_to_markdown_table(rows)
                if markdown_table:
                    elements.append(f"\n[TABLE_START]\n{markdown_table}\n[TABLE_END]\n")

        textboxes = _extract_docx_textboxes(document)
        if textboxes:
            elements.append("## Textboxes")
            elements.extend([f"- {textbox}" for textbox in textboxes])

        image_blobs = _collect_docx_image_blobs(document)
        if OCR_ENABLED and image_blobs:
            strategy = f"{strategy}|docx_images:{len(image_blobs)}"
            ocr_attempted_units = len(image_blobs)

            if Image is None:
                ocr_traces.append("docx_images:pillow_unavailable")
                strategy = f"{strategy}|ocr_dependency_missing:pillow"
            else:
                ocr_sections = []
                for image_index, blob in enumerate(image_blobs, start=1):
                    try:
                        with Image.open(io.BytesIO(blob)) as image_obj:
                            prepared = _prepare_image_for_ocr(image_obj.copy())
                    except Exception as exc:
                        ocr_traces.append(f"img{image_index}:decode_error:{exc.__class__.__name__}")
                        continue

                    ocr_result = _ocr_image_with_priority(prepared)
                    ocr_text_base = str(ocr_result.get("text") or "")
                    ocr_trace_base = str(ocr_result.get("trace") or "ocr:no_trace")
                    enhanced_text, layout_metadata = _augment_ocr_text_with_layout(
                        ocr_text_base,
                        list(ocr_result.get("blocks") or []),
                    )
                    if layout_metadata.get("table_added"):
                        ocr_layout_tables += 1
                    if layout_metadata.get("diagram_added"):
                        ocr_layout_diagrams += 1

                    ocr_trace = (
                        f"{ocr_trace_base}|layout_boxes:{int(layout_metadata.get('box_count', 0))}"
                        f"|layout_table:{int(bool(layout_metadata.get('table_added')))}"
                        f"|layout_diagram:{int(bool(layout_metadata.get('diagram_added')))}"
                    )
                    ocr_traces.append(f"img{image_index}:{ocr_trace}")

                    ocr_text = normalize_text(enhanced_text)
                    if not ocr_text:
                        continue

                    ocr_used_units += 1
                    ocr_sections.append(f"### Embedded image {image_index}\n{ocr_text}")

                if ocr_sections:
                    elements.append("## OCR embedded images")
                    elements.extend(ocr_sections)

            strategy = f"{strategy}|ocr_embedded_images:{ocr_used_units}/{len(image_blobs)}"
            strategy = (
                f"{strategy}|ocr_layout_tables:{ocr_layout_tables}|"
                f"ocr_layout_diagrams:{ocr_layout_diagrams}"
            )
            if ocr_traces:
                strategy = f"{strategy}|ocr_trace:{','.join(ocr_traces[:6])}"

        full_text = normalize_text("\n\n".join(elements))
        return _with_ocr_metadata(
            _response(full_text, None if full_text else _error("ERR_EMPTY", "Empty Docx"), strategy),
            attempted=ocr_attempted_units > 0,
            used=ocr_used_units > 0,
            trace=",".join(ocr_traces[:20]) if ocr_traces else None,
            pages=ocr_used_units,
            supplement_pages=ocr_used_units,
        )
    except Exception as exc:
        return _response("", _error("ERR_DOCX", str(exc)), strategy)


def _resolve_antiword_binary() -> str:
    local_candidate = Path(__file__).resolve().parent.parent / ".tools" / "antiword" / "bin" / "antiword"
    if local_candidate.exists() and os.access(local_candidate, os.X_OK):
        return str(local_candidate)

    if os.path.isabs(ANTIWORD_BIN) and os.path.exists(ANTIWORD_BIN):
        return ANTIWORD_BIN
    which_result = shutil.which(ANTIWORD_BIN)
    return which_result or ANTIWORD_BIN


def _resolve_antiword_share_dir() -> str:
    if ANTIWORD_SHARE_DIR and os.path.isdir(ANTIWORD_SHARE_DIR):
        return ANTIWORD_SHARE_DIR

    local_share_candidate = Path(__file__).resolve().parent.parent / ".tools" / "antiword" / "share" / "antiword"
    if local_share_candidate.is_dir():
        return str(local_share_candidate)
    return ""


def _extract_doc_binary_strings_fallback(content: bytes, max_lines: int = 1200) -> str:
    if not content:
        return ""

    lines: List[str] = []
    seen = set()
    for raw_chunk in re.findall(rb"[\x20-\x7E]{6,}", content):
        decoded = normalize_text(raw_chunk.decode("latin-1", errors="ignore"))
        if len(decoded) < 6:
            continue

        alpha_chars = sum(1 for char in decoded if char.isalpha())
        if alpha_chars < max(3, int(len(decoded) * 0.22)):
            continue

        if decoded in seen:
            continue
        seen.add(decoded)
        lines.append(decoded)

        if len(lines) >= max_lines:
            break

    return normalize_text("\n".join(lines))


def extract_text_from_doc_antiword(content: bytes) -> Dict[str, Any]:
    strategy = f"antiword:doc|limits:cpu={ANTIWORD_CPU_SECONDS},mem={ANTIWORD_MEMORY_MB}MB"
    if not content:
        return _response("", _error("ERR_EMPTY", "Empty"), strategy)

    antiword_binary = _resolve_antiword_binary()
    antiword_share_dir = _resolve_antiword_share_dir()
    strategy = f"{strategy}|bin:{os.path.basename(str(antiword_binary))}"
    if not antiword_binary or (not shutil.which(antiword_binary) and not os.path.exists(antiword_binary)):
        return _response(
            "",
            _error("ERR_ANTIWORD_BIN", f"Antiword binary not found ({ANTIWORD_BIN}). Configure ANTIWORD_BIN environment variable."),
            strategy,
        )

    try:
        with tempfile.TemporaryDirectory(prefix="parsing_doc_") as temp_dir:
            temp_path = os.path.join(temp_dir, "input.doc")
            with open(temp_path, "wb") as temp_file:
                temp_file.write(content)

            custom_env = os.environ.copy()
            if antiword_share_dir:
                custom_env["ANTIWORDHOME"] = antiword_share_dir

            command_specs = [
                ("utf8_mapping", [antiword_binary, "-m", "UTF-8.txt", "-w", "0", temp_path]),
                ("plain", [antiword_binary, "-w", "0", temp_path]),
            ]

            attempt_errors: List[str] = []

            for mode, command in command_specs:
                result = subprocess.run(
                    command,
                    capture_output=True,
                    text=True,
                    check=False,
                    env=custom_env,
                    timeout=PARSE_TIMEOUT_SECONDS,
                    preexec_fn=_antiword_preexec_limits(),
                )

                if result.returncode == 0 and result.stdout.strip():
                    mode_strategy = f"{strategy}|mode:{mode}"
                    if mode == "plain":
                        mode_strategy = f"{mode_strategy}|mapping_fallback"
                    return _response(normalize_text(result.stdout), None, mode_strategy)

                stderr_text = (result.stderr or "").strip()
                attempt_errors.append(f"{mode}:{stderr_text or f'code={result.returncode}'}")

                # Missing UTF mapping is common with system antiword installs.
                if mode == "utf8_mapping":
                    lowered = stderr_text.lower()
                    if "mapping file" in lowered or "utf-8.txt" in lowered:
                        continue

            fallback_text = _extract_doc_binary_strings_fallback(content)
            if fallback_text:
                fallback_strategy = f"{strategy}|mode:binary_strings_fallback"
                return _response(fallback_text, None, fallback_strategy)

            error_text = " | ".join(attempt_errors) if attempt_errors else "Antiword failed"
            return _response("", _error("ERR_ANTIWORD", error_text), strategy)
    except subprocess.TimeoutExpired:
        return _response("", _error("ERR_TIMEOUT", f"Antiword timed out after {PARSE_TIMEOUT_SECONDS}s"), strategy)
    except Exception as exc:
        return _response("", _error("ERR_ANTIWORD", str(exc)), strategy)


def extract_text_from_excel(content: bytes, extension: str) -> Dict[str, Any]:
    ext = _normalize_extension(extension) or "xlsx"
    engine = _excel_engine_for_extension(ext)
    excel_mode = "row_context+table" if EXCEL_INCLUDE_ROW_CONTEXT else "table_only"
    strategy = _with_lib_versions(f"pandas_excel:{ext}|engine:{engine or 'auto'}|mode:{excel_mode}", "pandas")
    if not content:
        return _response("", _error("ERR_EMPTY", "Empty"), strategy)

    try:
        started_at = time.time()

        try:
            workbook = pd.ExcelFile(io.BytesIO(content), engine=engine) if engine else pd.ExcelFile(io.BytesIO(content))
        except Exception:
            workbook = pd.ExcelFile(io.BytesIO(content))
            strategy = f"{strategy}|engine_fallback:auto"

        selected_sheets = workbook.sheet_names[:MAX_EXCEL_SHEETS]
        elements = []

        for sheet_name in selected_sheets:
            if _is_timed_out(started_at):
                return _response("", _error("ERR_TIMEOUT", f"Excel parsing exceeded {PARSE_TIMEOUT_SECONDS}s"), strategy)

            dataframe = workbook.parse(sheet_name=sheet_name, nrows=MAX_EXCEL_ROWS_PER_SHEET, dtype=str)
            headers, cleaned_dataframe = _clean_excel_headers_and_rows(dataframe)
            elements.append(f"## Sheet: {sheet_name}")

            preview_rows = [headers]

            for row_index, row_values in enumerate(cleaned_dataframe.itertuples(index=False, name=None)):
                normalized_values = [normalize_text("" if value is None else str(value)) for value in row_values]
                if row_index < EXCEL_TABLE_PREVIEW_ROWS:
                    preview_rows.append(normalized_values)

                if EXCEL_INCLUDE_ROW_CONTEXT:
                    row_context = ", ".join(
                        [
                            f"{str(header)}: {str(value)}"
                            for header, value in zip(headers, normalized_values)
                            if value
                        ]
                    )
                    if row_context:
                        elements.append(f"- {row_context}")

            markdown_preview = rows_to_markdown_table(preview_rows)
            if markdown_preview:
                elements.append(f"\n[TABLE_START]\n{markdown_preview}\n[TABLE_END]\n")

            if len(cleaned_dataframe) >= MAX_EXCEL_ROWS_PER_SHEET:
                elements.append(f"[WARNING] Row limit reached ({MAX_EXCEL_ROWS_PER_SHEET})")

        if len(workbook.sheet_names) > MAX_EXCEL_SHEETS:
            elements.append(f"[WARNING] Sheet limit reached ({MAX_EXCEL_SHEETS})")

        full_text = normalize_text("\n".join(elements))
        return _response(full_text, None if full_text else _error("ERR_EMPTY", "Empty Excel"), strategy)
    except Exception as exc:
        return _response("", _error("ERR_EXCEL", str(exc)), f"pandas_excel:{ext}")


def extract_text_from_pptx(content: bytes, extension: str) -> Dict[str, Any]:
    ext = _normalize_extension(extension) or "pptx"
    strategy = _with_lib_versions(f"python_pptx:{ext}", "python-pptx")
    if not content:
        return _response("", _error("ERR_EMPTY", "Empty"), strategy)

    try:
        started_at = time.time()
        presentation = Presentation(io.BytesIO(content))
        elements = []

        total_slides = len(presentation.slides)
        if total_slides > MAX_PPTX_SLIDES:
            elements.append(f"[WARNING] Slide limit reached ({MAX_PPTX_SLIDES}/{total_slides})")
            strategy = f"{strategy}|slide_limit:{MAX_PPTX_SLIDES}"

        def shape_sort_key(shape):
            return (
                _safe_int(getattr(shape, "top", 0)),
                _safe_int(getattr(shape, "left", 0)),
                _safe_int(getattr(shape, "shape_id", 0)),
            )

        def extract_chart_text(shape) -> List[str]:
            if not getattr(shape, "has_chart", False):
                return []

            chart = shape.chart
            chart_title = ""
            try:
                if chart.has_title and chart.chart_title and chart.chart_title.text_frame:
                    chart_title = normalize_text(chart.chart_title.text_frame.text)
            except Exception:
                chart_title = ""

            series_summaries = []
            for series in chart.series:
                series_name = normalize_text(str(getattr(series, "name", ""))) or "Series"
                try:
                    values = [str(value) for value in series.values]
                except Exception:
                    values = []

                if values:
                    preview = ", ".join(values[:12])
                    if len(values) > 12:
                        preview = f"{preview}, ..."
                    series_summaries.append(f"{series_name}: {preview}")

            if not series_summaries and chart_title:
                return [f"Chart: {chart_title}"]
            if series_summaries and chart_title:
                return [f"Chart: {chart_title} | {' ; '.join(series_summaries)}"]
            if series_summaries:
                return [f"Chart: {' ; '.join(series_summaries)}"]
            return ["Chart"]

        def extract_shape_texts(shape) -> Iterable[str]:
            if getattr(shape, "has_text_frame", False) and shape.text:
                text = normalize_text(shape.text)
                if text:
                    yield text

            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    rows.append([normalize_text(cell.text) for cell in row.cells])
                markdown_table = rows_to_markdown_table(rows)
                if markdown_table:
                    yield f"[TABLE_START]\n{markdown_table}\n[TABLE_END]"

            for chart_text in extract_chart_text(shape):
                normalized_chart = normalize_text(chart_text)
                if normalized_chart:
                    yield normalized_chart

            if hasattr(shape, "shapes"):
                for nested_shape in shape.shapes:
                    yield from extract_shape_texts(nested_shape)

        for slide_num, slide in enumerate(presentation.slides, start=1):
            if _is_timed_out(started_at):
                return _response("", _error("ERR_TIMEOUT", f"PPTX parsing exceeded {PARSE_TIMEOUT_SECONDS}s"), strategy)
            if slide_num > MAX_PPTX_SLIDES:
                break

            layout_name = normalize_text(getattr(slide.slide_layout, "name", ""))
            slide_title = ""
            try:
                if slide.shapes.title and slide.shapes.title.text:
                    slide_title = normalize_text(slide.shapes.title.text)
            except Exception:
                slide_title = ""

            slide_header = f"## Slide {slide_num}"
            if layout_name:
                slide_header = f"{slide_header} | Layout: {layout_name}"
            elements.append(slide_header)
            if slide_title:
                elements.append(f"### Slide Title: {slide_title}")

            for shape in sorted(list(slide.shapes), key=shape_sort_key):
                elements.extend([text for text in extract_shape_texts(shape) if text])

            if slide.has_notes_slide:
                notes_text = normalize_text(slide.notes_slide.notes_text_frame.text)
                if notes_text:
                    elements.append(f"## Slide {slide_num} Notes\n{notes_text}")

        full_text = normalize_text("\n\n".join(elements))
        return _response(full_text, None if full_text else _error("ERR_EMPTY", "Empty"), strategy)
    except Exception as exc:
        return _response("", _error("ERR_PPTX", str(exc)), strategy)


def extract_text_locally(content: bytes, extension: str) -> Dict[str, Any]:
    start_time = time.time()
    ext = _normalize_extension(extension)

    if not content:
        result = _response("", _error("ERR_EMPTY", "No content"), f"unhandled:{ext or 'none'}")
        result["parse_time_seconds"] = float(round(time.time() - start_time, 4))
        return result

    if len(content) > MAX_DOCUMENT_BYTES:
        result = _response(
            "",
            _error("ERR_SIZE_LIMIT", f"Document size {len(content)} exceeds limit {MAX_DOCUMENT_BYTES} bytes"),
            f"guard:size_limit:{ext or 'unknown'}",
        )
        result["parse_time_seconds"] = float(round(time.time() - start_time, 4))
        return result

    detected_format = _detect_binary_format(content)
    if _is_mime_mismatch_for_textlike(ext, detected_format):
        result = _response(
            "",
            _error("ERR_MIME_MISMATCH", f"Extension '{ext}' does not match binary signature '{detected_format}'"),
            f"guard:mime_mismatch:{ext}|signature:{detected_format}",
        )
        result["parse_time_seconds"] = float(round(time.time() - start_time, 4))
        return result

    extension_format_mismatch = (
        (ext == "pdf" and detected_format not in {"pdf", "unknown"})
        or (ext in {"doc", "xls", "ppt"} and detected_format not in {"ole2", "unknown"})
        or (ext in {"docx", "xlsx", "pptx", "xlsm", "xlsb"} and detected_format not in {"ooxml", "unknown"})
        or (ext in _IMAGE_EXTENSIONS and detected_format not in {"jpeg", "png", "tiff", "webp", "unknown"})
    )
    format_suffix = f"|signature:{detected_format}" if extension_format_mismatch else ""

    if ext == "txt":
        text = normalize_text(safe_decode(content))
        result = _response(text, None if text else _error("ERR_EMPTY", "Empty TXT"), f"txt_decoder{format_suffix}")
    elif ext == "pdf":
        result = extract_text_from_pdf_pdfplumber(content, ext)
        if extension_format_mismatch:
            result["parser_strategy"] = f"{result.get('parser_strategy', 'pdfplumber:pdf')}{format_suffix}"
    elif ext in {"docx", "docm"}:
        result = extract_text_from_docx(content, ext)
        if extension_format_mismatch:
            result["parser_strategy"] = f"{result.get('parser_strategy', 'python_docx:docx')}{format_suffix}"
    elif ext == "doc":
        result = extract_text_from_doc_antiword(content)
        if extension_format_mismatch:
            result["parser_strategy"] = f"{result.get('parser_strategy', 'antiword:doc')}{format_suffix}"
    elif ext in {"xls", "xlsx", "xlsm", "xlsb"}:
        result = extract_text_from_excel(content, ext)
        if extension_format_mismatch:
            result["parser_strategy"] = f"{result.get('parser_strategy', 'pandas_excel_row_context')}{format_suffix}"
    elif ext in {"pptx"}:
        result = extract_text_from_pptx(content, ext)
        if extension_format_mismatch:
            result["parser_strategy"] = f"{result.get('parser_strategy', 'python_pptx:pptx')}{format_suffix}"
    elif ext in _IMAGE_EXTENSIONS:
        result = extract_text_from_image_ocr(content, ext)
        if extension_format_mismatch:
            result["parser_strategy"] = f"{result.get('parser_strategy', 'ocr_image')}{format_suffix}"
    elif ext in TEXTLIKE_EXTENSIONS:
        result = extract_text_from_textlike(content, ext)
    elif detected_format in {"jpeg", "png", "tiff", "webp"}:
        result = extract_text_from_image_ocr(content, detected_format)
        result["parser_strategy"] = f"{result.get('parser_strategy', 'ocr_image')}|detected_format_routing"
    else:
        result = _response("", _error("ERR_UNHANDLED_EXTENSION", f"No local parser available for {ext}"), f"unhandled:{ext}")

    if len(result.get("text") or "") > MAX_CHUNK_CHARS * 200:
        result["parser_strategy"] = f"{result.get('parser_strategy', 'local')}|very_large_text"

    result["parse_time_seconds"] = float(round(time.time() - start_time, 4))
    return result


def extract_text_locally_with_retry(
    content: bytes,
    extension: str,
    max_retries: int = LOCAL_PARSE_MAX_RETRIES,
) -> Dict[str, Any]:
    started_at = time.time()
    attempts = max(1, int(max_retries))
    result: Dict[str, Any] = {}

    for attempt in range(1, attempts + 1):
        result = extract_text_locally(content, extension)
        parser_error = result.get("parser_error")
        if not parser_error:
            break
        if attempt >= attempts or not is_retryable_error(parser_error):
            break

    base_strategy = result.get("parser_strategy", "local")
    result["parser_strategy"] = f"{base_strategy}|attempts:{attempt}"
    result["parse_time_seconds"] = float(round(time.time() - started_at, 4))
    return result
