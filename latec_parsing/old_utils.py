IDentifie les défaut de parsing de ceci (pas les défaut liés au nombre de tokens pas exactes etc les vraies défaut qui vont faire du mauvais parsing, par exemple les unnames qui se répète X fois dans les excel ou autre 

import io
import os
import json
import re
import tempfile
import time
import zipfile 
import subprocess
import uuid
from typing import Any, Dict, List

import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
from lxml import etree 

from langchain_text_splitters import RecursiveCharacterTextSplitter, MarkdownHeaderTextSplitter

from pyspark.sql import functions as F
from pyspark.sql import types as T
from pyspark.sql.functions import pandas_udf

# ============================================================
# Constants & Configs
# ============================================================
TOKENIZER_ENCODING = "cl100k_base"

# TODO: Replace antiword with LibreOffice headless in the future
# Command would be: soffice --headless --convert-to docx /tmp/input.doc --outdir /tmp/
ANTIWORD_BIN = "/Workspace/Users/jules.gourio.external@latecoere.aero/Parsing_RAG/antiword_local/usr/bin/antiword"
ANTIWORD_SHARE_DIR = "/Workspace/Users/jules.gourio.external@latecoere.aero/Parsing_RAG/antiword_local/usr/share/antiword"

TEXTLIKE_EXTENSIONS = {"html", "xml", "md", "csv", "json", "tsv"}

# ============================================================
# Schemas
# ============================================================
TEXT_PARSE_SCHEMA = T.StructType([
    T.StructField("text", T.StringType(), True),
    T.StructField("parser_error", T.StringType(), True),
    T.StructField("parser_strategy", T.StringType(), True),
    T.StructField("parse_time_seconds", T.FloatType(), True), 
])

CHUNK_SCHEMA = T.ArrayType(T.StructType([
    T.StructField("chunk_index", T.IntegerType(), True),
    T.StructField("chunk_text", T.StringType(), True),
    T.StructField("chunk_char_count", T.IntegerType(), True),
    T.StructField("chunk_token_count", T.IntegerType(), True),
    T.StructField("chunk_content_type", T.StringType(), True),
    T.StructField("metadata", T.MapType(T.StringType(), T.StringType()), True)
]))

# ============================================================
# Basic Helpers
# ============================================================
def safe_decode(raw_bytes: bytes, encodings: List[str] = None) -> str:
    if not raw_bytes: return ""
    for enc in (encodings or ["utf-8", "utf-8-sig", "cp1252", "latin-1"]):
        try: return raw_bytes.decode(enc)
        except: continue
    return raw_bytes.decode("utf-8", errors="ignore")

def strip_html_tags(text: str) -> str:
    """Robust HTML cleaning using BeautifulSoup instead of volatile regex."""
    if not text: return ""
    soup = BeautifulSoup(text, "html.parser")
    # Remove script and style elements completely
    for script in soup(["script", "style"]):
        script.extract()
    return soup.get_text(separator=" ", strip=True)

def normalize_text(text: str) -> str:
    if not text: return ""
    text = text.replace("\r\n", "\n").replace("\r", "\n").replace("\x00", "")
    text = re.sub(r"[ \t]+", " ", text)
    return re.sub(r"\n{3,}", "\n\n", text).strip()

def rows_to_markdown_table(rows: List[List[Any]]) -> str:
    """Standard markdown table generation (Keeps the |--|--| line for LLM comprehension)"""
    if not rows: return ""
    normalized_rows = [[normalize_text(str(cell) if pd.notna(cell) and cell is not None else "") for cell in row] for row in rows]
    col_count = max((len(row) for row in normalized_rows), default=0)
    if col_count == 0: return ""
    normalized_rows = [row + [""] * (col_count - len(row)) for row in normalized_rows]
    header = normalized_rows[0]
    body = normalized_rows[1:] if len(normalized_rows) > 1 else []

    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * col_count) + " |"]
    for row in body: lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)

def get_tiktoken_encoder():
    """Instantiated locally to prevent Spark Pickling errors on globals."""
    import tiktoken
    return tiktoken.get_encoding(TOKENIZER_ENCODING)

def count_tokens(text: str) -> int:
    if not text: return 0
    try:
        enc = get_tiktoken_encoder()
        return len(enc.encode(text))
    except:
        return max(1, len(text) // 4)

# ============================================================
# Extraction Functions
# ============================================================
def extract_text_from_textlike(content: bytes, extension: str) -> Dict[str, str]:
    ext = (extension or "").lower()
    if not content: return {"text": "", "parser_error": "Empty", "parser_strategy": f"textlike:{ext}"}
    
    decoded = safe_decode(content)
    if ext in {"html", "htm"}: 
        decoded = strip_html_tags(decoded)
    elif ext in {"csv", "tsv"}:
        # Handle CSVs like Excel to avoid giant markdown tables
        try:
            df = pd.read_csv(io.StringIO(decoded), sep="\t" if ext == "tsv" else ",")
            elements = []
            headers = df.columns.tolist()
            for idx, row in df.iterrows():
                row_context = ", ".join([f"{str(h)}: {str(row[h])}" for h in headers if pd.notna(row[h])])
                elements.append(f"- {row_context}")
            decoded = "\n".join(elements)
        except Exception:
            decoded = f"### CSV Data\n[TABLE_START]\n{decoded}\n[TABLE_END]"
            
    decoded = normalize_text(decoded)
    return {"text": decoded, "parser_error": None if decoded else "Empty decoded", "parser_strategy": f"textlike:{ext}"}

def extract_text_from_pdf_pdfplumber(content: bytes, extension: str) -> Dict[str, str]:
    """Uses pdfplumber for better multi-column layout detection and table extraction."""
    ext = (extension or "pdf").lower()
    if not content: return {"text": "", "parser_error": "Empty", "parser_strategy": f"pdfplumber:{ext}"}
    
    full_text = ""
    try:
        elements = []
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page in pdf.pages:
                # 1. Extract text keeping reading layout (fixes double column mixing)
                text = page.extract_text(layout=True)
                if text:
                    elements.append(text)
                
                # 2. Extract tables explicitly
                tables = page.extract_tables()
                for table in tables:
                    md_table = rows_to_markdown_table(table)
                    if md_table:
                        elements.append(f"\n[TABLE_START]\n{md_table}\n[TABLE_END]\n")
                        
        full_text = normalize_text("\n\n".join(elements))
        return {"text": full_text, "parser_error": None if full_text else "Empty PDF", "parser_strategy": f"pdfplumber:{ext}"}
    except Exception as e:
        return {"text": "", "parser_error": str(e), "parser_strategy": f"pdfplumber:{ext}"}

def extract_text_from_docx(content: bytes, extension: str) -> Dict[str, str]:
    ext = (extension or "docx").lower()
    try:
        doc = Document(io.BytesIO(content))
        elements = []
        for child in doc.element.body.iterchildren():
            if child.tag.endswith('p'):
                from docx.text.paragraph import Paragraph
                para = Paragraph(child, doc)
                style_name = para.style.name.lower()
                txt = normalize_text(para.text)
                if txt:
                    # Inject markdown headers to preserve hierarchy for Langchain
                    if "heading 1" in style_name: txt = f"# {txt}"
                    elif "heading 2" in style_name: txt = f"## {txt}"
                    elif "heading 3" in style_name: txt = f"### {txt}"
                    elements.append(txt)
            elif child.tag.endswith('tbl'):
                from docx.table import Table
                rows = [[cell.text.strip() for cell in row.cells] for row in Table(child, doc).rows]
                md_table = rows_to_markdown_table(rows)
                if md_table: elements.append(f"\n[TABLE_START]\n{md_table}\n[TABLE_END]\n")
                
        full_text = normalize_text("\n\n".join(elements))
        return {"text": full_text, "parser_error": None if full_text else "Empty Docx", "parser_strategy": f"python_docx:{ext}"}
    except Exception as e:
        return {"text": "", "parser_error": str(e), "parser_strategy": f"python_docx:{ext}"}

def extract_text_from_doc_antiword(content: bytes) -> Dict[str, str]:
    """Legacy .doc parsing. TODO: Upgrade to LibreOffice headless (soffice)"""
    if not content: return {"text": "", "parser_error": "Empty", "parser_strategy": "antiword:doc"}

    unique_id = uuid.uuid4().hex[:8]
    temp_path = f"/tmp/antiword_parse_{unique_id}.doc"
    try:
        with open(temp_path, "wb") as f:
            f.write(content)
        custom_env = os.environ.copy()
        custom_env["ANTIWORDHOME"] = ANTIWORD_SHARE_DIR
        result = subprocess.run([ANTIWORD_BIN, "-m", "UTF-8.txt", "-w", "0", temp_path],
                                capture_output=True, text=True, check=False, env=custom_env)

        if result.returncode == 0 and result.stdout.strip():
            full_text = normalize_text(result.stdout)
            return {"text": full_text, "parser_error": None, "parser_strategy": "antiword:doc"}
        else:
            return {"text": "", "parser_error": result.stderr.strip() or "Antiword failed", "parser_strategy": "antiword:doc"}
    except Exception as e:
        return {"text": "", "parser_error": str(e), "parser_strategy": "antiword:doc"}
    finally:
        if os.path.exists(temp_path):
            try: os.remove(temp_path)
            except: pass

def extract_text_from_excel(content: bytes, extension: str) -> Dict[str, str]:
    """Extracts Excel files row-by-row to prevent giant unbreakable Markdown tables."""
    ext = (extension or "xlsx").lower()
    try:
        dfs = pd.read_excel(io.BytesIO(content), sheet_name=None)
        elements = []
        for sheet_name, df in dfs.items():
            elements.append(f"## Sheet: {sheet_name}") # Preserves hierarchy
            headers = df.columns.tolist()
            for idx, row in df.iterrows():
                # Creates a contextual sentence per row
                row_context = ", ".join([f"{str(h)}: {str(row[h])}" for h in headers if pd.notna(row[h])])
                elements.append(f"- {row_context}")
        
        full_text = normalize_text("\n".join(elements))
        return {"text": full_text, "parser_error": None if full_text else "Empty Excel", "parser_strategy": f"pandas_excel_row_context:{ext}"}
    except Exception as e:
        return {"text": "", "parser_error": str(e), "parser_strategy": f"pandas_excel:{ext}"}

def extract_text_from_pptx(content: bytes, extension: str) -> Dict[str, str]:
    ext = (extension or "pptx").lower()
    try:
        prs = Presentation(io.BytesIO(content))
        elements = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            elements.append(f"## Slide {slide_num}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    elements.append(normalize_text(shape.text))
        full_text = normalize_text("\n\n".join(elements))
        return {"text": full_text, "parser_error": None if full_text else "Empty", "parser_strategy": f"python_pptx:{ext}"}
    except Exception as e:
        return {"text": "", "parser_error": str(e), "parser_strategy": f"python_pptx:{ext}"}

# ============================================================
# Local Extraction Wrapper
# ============================================================
def extract_text_locally(content: bytes, extension: str) -> Dict[str, Any]:
    start_time = time.time()
    ext = (extension or "").lower()

    if ext == "txt":
        text = normalize_text(safe_decode(content))
        res = {"text": text, "parser_error": None if text else "Empty TXT", "parser_strategy": "txt_decoder"}
    elif ext == "pdf":
        res = extract_text_from_pdf_pdfplumber(content, ext)
    elif ext in {"docx", "docm"}:
        res = extract_text_from_docx(content, ext)
    elif ext == "doc":
        res = extract_text_from_doc_antiword(content)
    elif ext in {"xls", "xlsx"}:
        res = extract_text_from_excel(content, ext)
    elif ext in {"pptx"}:
        res = extract_text_from_pptx(content, ext)
    elif ext in TEXTLIKE_EXTENSIONS:
        res = extract_text_from_textlike(content, ext)
    else:
        res = {"text": "", "parser_error": f"No local parser available for {ext}", "parser_strategy": f"unhandled:{ext}"}

    res["parse_time_seconds"] = float(round(time.time() - start_time, 2))
    return res

# ============================================================
# Chunking Engine (Semantic + Token-aware)
# ============================================================
def split_text_to_chunks(text: str, chunk_size_tokens: int, chunk_overlap_tokens: int) -> List[Dict[str, Any]]:
    text = normalize_text(text)
    if not text: return []

    # 1. Semantic Splitting: Keep headers attached to their content
    headers_to_split_on = [
        ("#", "Header 1"),
        ("##", "Header 2"),
        ("###", "Header 3"),
    ]
    markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
    semantic_splits = markdown_splitter.split_text(text)

    # 2. Token-aware Recursive Splitting (Exact token count, respects tables better)
    # We add [TABLE_START] and [TABLE_END] to separators so it avoids breaking tables if possible
    custom_separators = [
        "\n\n", 
        "\n[TABLE_START]\n", 
        "\n[TABLE_END]\n", 
        "\n", 
        ". ", 
        " ", 
        ""
    ]
    
    token_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
        encoding_name=TOKENIZER_ENCODING,
        chunk_size=chunk_size_tokens,
        chunk_overlap=chunk_overlap_tokens,
        separators=custom_separators
    )

    final_chunks = []
    chunk_index = 0

    for doc in semantic_splits:
        # doc.page_content holds the text, doc.metadata holds the headers (e.g., {"Header 1": "Introduction"})
        sub_chunks = token_splitter.split_text(doc.page_content)
        
        for sub_chunk in sub_chunks:
            # Re-inject headers into the text for the LLM to see context
            header_context = " > ".join(doc.metadata.values())
            enriched_text = f"[{header_context}]\n{sub_chunk}" if header_context else sub_chunk
            
            # Simple AI guardrail: Skip identical consecutive chunks
            if final_chunks and final_chunks[-1]["chunk_text"] == enriched_text:
                continue

            final_chunks.append({
                "chunk_index": chunk_index,
                "chunk_text": enriched_text,
                "chunk_char_count": len(enriched_text),
                "chunk_token_count": count_tokens(enriched_text),
                "chunk_content_type": "table" if "[TABLE_START]" in enriched_text else "text",
                "metadata": doc.metadata
            })
            chunk_index += 1

    return final_chunks

# ============================================================
# PySpark Pandas UDFs (Apache Arrow)
# ============================================================

@pandas_udf(TEXT_PARSE_SCHEMA)
def extract_local_text_udf(content_series: pd.Series, ext_series: pd.Series) -> pd.DataFrame:
    results = []
    for content, ext in zip(content_series, ext_series):
        if content is None:
            results.append({"text": "", "parser_error": "No Content", "parser_strategy": "none", "parse_time_seconds": 0.0})
        else:
            results.append(extract_text_locally(content, str(ext)))
    return pd.DataFrame(results)

@pandas_udf(CHUNK_SCHEMA)
def build_chunks_udf(text_series: pd.Series, sz_series: pd.Series, ov_series: pd.Series) -> pd.Series:
    results = []
    for text, sz, ov in zip(text_series, sz_series, ov_series):
        if not text:
            results.append([])
        else:
            # extension is no longer needed since markdown splitting handles docx/excel logic internally
            results.append(split_text_to_chunks(str(text), int(sz), int(ov)))
    return pd.Series(results)

@pandas_udf(T.IntegerType())
def token_count_udf(text_series: pd.Series) -> pd.Series:
    results = [count_tokens(str(text)) if text else 0 for text in text_series]
    return pd.Series(results)