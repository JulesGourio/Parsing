import gc
import json
import logging
import mimetypes
import os
import re
import time
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from dataclasses import replace

import fitz  # PyMuPDF


# ------------------------------------------------------------------------------
# Configuration
# ------------------------------------------------------------------------------
DEFAULT_CONTEXT_WINDOW = 250
OFFLINE_MODELS_DIR = Path(__file__).resolve().parent / "docling_models"
SCRIPT_DIR = Path(__file__).resolve().parent
# INPUT_SOURCE = 

OUTPUT_MD = SCRIPT_DIR / "benchmark_nextgen_report.md"
RUNS_PER_METHOD = 1
INCLUDE_RAW_MARKDOWN = True
INCLUDE_IMAGE_CONTEXTS = False
INCLUDE_PDF_GALLERY = True
GALLERY_MAX_IMAGES = 500
# Activer uniquement si des modeles OCR sont disponibles dans docling_models/
FORCE_FULL_PAGE_OCR = False
IMAGES_PER_PDF = 500          # Nombre max d'images extraites via PyMuPDF par document
FIX_GARBLED_TEXT = True      # Correction ASCII+29 ligne par ligne apres pypdfium2 (mixte-encodage)

os.environ["HF_HUB_OFFLINE"] = "1"
os.environ["TRANSFORMERS_OFFLINE"] = "1"
os.environ["HF_HUB_DISABLE_TELEMETRY"] = "1"

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S",
)
logger = logging.getLogger(__name__)
logging.getLogger("docling").setLevel(logging.WARNING)
logging.getLogger("RapidOCR").setLevel(logging.ERROR)


# ------------------------------------------------------------------------------
# Imports Docling
# ------------------------------------------------------------------------------
from docling.document_converter import DocumentConverter

try:
    from docling.document_converter import PdfFormatOption
except Exception:
    PdfFormatOption = None  # type: ignore[assignment]

try:
    from docling.datamodel.base_models import InputFormat
except Exception:
    InputFormat = None  # type: ignore[assignment]

try:
    from docling.datamodel.pipeline_options import (
        PdfPipelineOptions,
    )
except Exception:
    PdfPipelineOptions = None  # type: ignore[assignment]

try:
    from docling.pipeline.vlm_pipeline import VlmPipeline
    from docling.datamodel.pipeline_options import VlmPipelineOptions

    try:
        from docling.datamodel.pipeline_options import VlmConvertOptions

        HAS_VLM_PRESET = True
    except Exception:
        from docling.datamodel import vlm_model_specs

        HAS_VLM_PRESET = False
    VLM_AVAILABLE = True
except Exception:
    logger.warning("Pipeline VLM non detecte. Mettez a jour docling.")
    VLM_AVAILABLE = False

# ------------------------------------------------------------------------------
# Modeles internes
# ------------------------------------------------------------------------------
@dataclass
class ImageContext:
    image_index: int
    image_ref: str
    context_before: str
    context_after: str
    llm_prompt: str
    simulated_summary: str


@dataclass
class ParsedArtifact:
    source: str
    source_type: str
    method: str
    parse_elapsed_s: float
    markdown: str
    stats: Dict[str, int]
    image_contexts: List[ImageContext]
    warnings: List[str]
    error: Optional[str] = None


@dataclass
class MethodBenchmark:
    method: str
    label: str
    runs: List[ParsedArtifact]

    @property
    def best(self) -> ParsedArtifact:
        valid = [r for r in self.runs if not r.error]
        if valid:
            return min(valid, key=lambda r: r.parse_elapsed_s)
        return self.runs[-1]

    @property
    def avg_parse_s(self) -> float:
        vals = [r.parse_elapsed_s for r in self.runs]
        return sum(vals) / max(1, len(vals))


SUPPORTED_METHODS = (
    "docling_standard",
    "fitz_direct",
    "docling_vlm_granite",
)

METHOD_LABELS = {
    "docling_standard": "Docling Standard (pypdfium2)",
    "fitz_direct": "PyMuPDF fitz natif",
    "docling_vlm_granite": "Docling VLM Granite",
}


# ------------------------------------------------------------------------------
# Utilitaires
# ------------------------------------------------------------------------------
def free_memory() -> None:
    gc.collect()


def is_url(source: str) -> bool:
    return source.startswith("http://") or source.startswith("https://")


def detect_source_type(source: str) -> str:
    if is_url(source):
        return "url"
    p = Path(source)
    if p.is_dir():
        return "directory"
    if p.is_file():
        return "file"
    return "unknown"


def sanitize_base64_for_report(md_text: str) -> str:
    return re.sub(
        r"data:image/[a-zA-Z]+;base64,[a-zA-Z0-9+/=]+",
        "[BASE64_IMAGE_DATA_EXTRACTED]",
        md_text,
    )


def analyze_markdown_content(md_text: str) -> Dict[str, int]:
    table_regex = re.compile(r"(?m)^\|.*\|\n^\|(?:[-:]+[-:| ]+)+\|\n(?:^\|.*\|\n)*")
    return {
        "chars": len(md_text),
        "words": len(re.findall(r"\b\w+\b", md_text)),
        "math_inline": len(re.findall(r"(?<!\$)\$(?!\$)(.*?)(?<!\$)\$(?!\$)", md_text)),
        "math_block": len(re.findall(r"\$\$(.*?)\$\$", md_text, re.DOTALL)),
        "tables_count": len(table_regex.findall(md_text)),
        "images_refs": count_image_refs(md_text),
    }


def count_image_refs(md_text: str) -> int:
    return len(find_image_markers(md_text))


def find_image_markers(md_text: str) -> List[Tuple[int, int, str]]:
    # Couvre les sorties frequentes Docling/Markdown/HTML pour eviter les faux 0.
    patterns = [
        re.compile(r"!\[[^\]]*\]\(([^\)]+)\)"),
        re.compile(r"<img\s+[^>]*src=[\"']([^\"']+)[\"']", re.IGNORECASE),
        re.compile(r"^\s*\[[^\]]+\]:\s*(\S+\.(?:png|jpg|jpeg|webp|gif|bmp|tiff?))(?:\s+.*)?$", re.IGNORECASE | re.MULTILINE),
        re.compile(r"\b(https?://\S+?\.(?:png|jpg|jpeg|webp|gif|bmp|tiff?))\b", re.IGNORECASE),
        re.compile(r"\[BASE64_IMAGE_DATA_EXTRACTED\]"),
        re.compile(r"<!--\s*image\s*-->", re.IGNORECASE),
        re.compile(r"<!--\s*figure\s*-->", re.IGNORECASE),
    ]

    markers: List[Tuple[int, int, str]] = []
    for pattern in patterns:
        for m in pattern.finditer(md_text):
            ref = m.group(1) if m.groups() else "docling_image_marker"
            markers.append((m.start(), m.end(), ref))

    # Dedup des doublons exacts (meme offset + meme ref)
    seen = set()
    deduped: List[Tuple[int, int, str]] = []
    for marker in sorted(markers, key=lambda x: (x[0], x[1], x[2])):
        key = (marker[0], marker[1], marker[2])
        if key in seen:
            continue
        seen.add(key)
        deduped.append(marker)
    return deduped


def truncate_for_display(text: str, max_len: int = 500) -> str:
    cleaned = re.sub(r"\s+", " ", text).strip()
    if len(cleaned) <= max_len:
        return cleaned
    return cleaned[: max_len - 3] + "..."


def extract_image_contexts(md_text: str, context_window: int = DEFAULT_CONTEXT_WINDOW) -> List[ImageContext]:
    matches = sorted(find_image_markers(md_text), key=lambda item: item[0])

    contexts: List[ImageContext] = []
    for idx, (start, end, image_ref) in enumerate(matches, start=1):
        before = md_text[max(0, start - context_window) : start]
        after = md_text[end : min(len(md_text), end + context_window)]
        prompt = (
            "Tu es un assistant de resume d'image pour un pipeline RAG. "
            "N'utilise pas d'OCR. Deduis uniquement un resume semantique probable de l'image "
            "a partir du contexte textuel local.\n\n"
            f"Image ref: {image_ref}\n"
            f"Contexte avant ({len(before)} chars):\n{before}\n\n"
            f"Contexte apres ({len(after)} chars):\n{after}\n\n"
            "Retour attendu: 3 a 5 phrases, ton neutre, focus sur informations utiles pour retrieval."
        )
        simulated = (
            "[SIMULATION - NON ENVOYE AU LLM] Resume image probable base sur contexte local. "
            "Cette section sera remplacee par la sortie d'un LLM dans l'integration finale."
        )
        contexts.append(
            ImageContext(
                image_index=idx,
                image_ref=image_ref,
                context_before=before,
                context_after=after,
                llm_prompt=prompt,
                simulated_summary=simulated,
            )
        )

    return contexts


def extract_pdf_images_to_dir(
    pdf_path: Path,
    images_dir: Path,
    max_images: int = 8,
    min_area_ratio: float = 0.02,
) -> List[Path]:
    images_dir.mkdir(parents=True, exist_ok=True)
    saved_images: List[Path] = []

    # Keep images that occupy significant area on page (more likely real figures/charts).
    candidates: List[Tuple[float, int, int, int]] = []  # (area_ratio, xref, width, height)

    with fitz.open(str(pdf_path)) as doc:
        seen_xrefs = set()
        for page in doc:
            page_area = float(page.rect.width * page.rect.height) if page.rect else 0.0
            if page_area <= 0:
                continue

            for image_info in page.get_images(full=True):
                xref = image_info[0]
                if xref in seen_xrefs:
                    continue
                width = int(image_info[2]) if len(image_info) > 2 else 0
                height = int(image_info[3]) if len(image_info) > 3 else 0

                try:
                    rects = page.get_image_rects(xref)
                except Exception:
                    continue

                if not rects:
                    continue

                best_ratio = 0.0
                for rect in rects:
                    ratio = float((rect.width * rect.height) / page_area)
                    if ratio > best_ratio:
                        best_ratio = ratio

                if best_ratio < min_area_ratio:
                    continue

                seen_xrefs.add(xref)
                candidates.append((best_ratio, xref, width, height))

        candidates.sort(key=lambda x: x[0], reverse=True)
        for _, xref, width, height in candidates[:max_images]:
            try:
                extracted = doc.extract_image(xref)
            except Exception:
                continue

            img_bytes = extracted.get("image")
            img_ext = extracted.get("ext", "png")
            if not img_bytes:
                continue

            out_name = f"{pdf_path.stem}_xref{xref}_{width}x{height}.{img_ext}"
            out_path = images_dir / out_name
            try:
                out_path.write_bytes(img_bytes)
            except Exception:
                continue

            saved_images.append(out_path)

    return saved_images


def extract_figure_crops_to_dir(
    pdf_path: Path,
    images_dir: Path,
    max_images: int = 8,
) -> List[Path]:
    """Crop les images embarquees dans le PDF autour de leurs vraies coordonnees."""
    images_dir.mkdir(parents=True, exist_ok=True)
    saved_images: List[Path] = []

    with fitz.open(str(pdf_path)) as doc:
        seen_xrefs: set = set()
        for page_idx, page in enumerate(doc, start=1):
            if len(saved_images) >= max_images:
                break

            page_area = float(page.rect.width * page.rect.height)
            if page_area <= 0:
                continue

            for img_info in page.get_images(full=True):
                if len(saved_images) >= max_images:
                    break
                xref = img_info[0]
                if xref in seen_xrefs:
                    continue

                try:
                    rects = page.get_image_rects(xref)
                except Exception:
                    continue
                if not rects:
                    continue

                # Choisir le rect le plus grand
                rect = max(rects, key=lambda r: r.width * r.height)
                ratio = float(rect.width * rect.height) / page_area
                if ratio < 0.015:
                    continue

                seen_xrefs.add(xref)

                # Crop autour des vraies coordonnees + petite marge
                pad = 4.0
                clip = fitz.Rect(
                    max(0.0, rect.x0 - pad),
                    max(0.0, rect.y0 - pad),
                    min(page.rect.width, rect.x1 + pad),
                    min(page.rect.height, rect.y1 + pad),
                )
                if clip.width < 60 or clip.height < 60:
                    continue

                try:
                    pix = page.get_pixmap(clip=clip, matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                except Exception:
                    continue

                out_name = f"{pdf_path.stem}_fig_p{page_idx:03d}_x{xref}.png"
                out_path = images_dir / out_name
                try:
                    pix.save(str(out_path))
                except Exception:
                    continue

                saved_images.append(out_path)

    return saved_images


def extract_pptx_images_to_dir(
    pptx_path: Path,
    images_dir: Path,
    max_images: int = 30,
) -> List[Path]:
    """Extrait les images embarquees dans un fichier PPTX via python-pptx."""
    images_dir.mkdir(parents=True, exist_ok=True)
    saved: List[Path] = []
    try:
        from pptx import Presentation  # type: ignore[import]
        prs = Presentation(str(pptx_path))
        idx = 0
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if len(saved) >= max_images:
                    break
                # shape_type 13 = PICTURE
                if not hasattr(shape, "image"):
                    continue
                try:
                    image = shape.image
                    ext = image.ext or "png"
                    out_path = images_dir / f"{pptx_path.stem}_s{slide_num:03d}_i{idx:03d}.{ext}"
                    out_path.write_bytes(image.blob)
                    saved.append(out_path)
                    idx += 1
                except Exception:
                    continue
    except ImportError:
        logger.warning("python-pptx non disponible — pip install python-pptx pour extraire les images PPTX")
    except Exception as e:
        logger.warning("Erreur extraction images PPTX: %s", e)
    logger.info("%d image(s) extraite(s) depuis PPTX %s", len(saved), pptx_path.name)
    return saved


def extract_docx_images_to_dir(
    docx_path: Path,
    images_dir: Path,
    max_images: int = 30,
) -> List[Path]:
    """Extrait les images embarquees dans un fichier DOCX via zipfile (word/media/)."""
    import zipfile
    images_dir.mkdir(parents=True, exist_ok=True)
    saved: List[Path] = []
    valid_ext = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp"}
    try:
        with zipfile.ZipFile(str(docx_path), "r") as z:
            media = [f for f in z.namelist() if f.startswith("word/media/")]
            for entry in media[:max_images]:
                if Path(entry).suffix.lower() not in valid_ext:
                    continue
                out_path = images_dir / f"{docx_path.stem}_{Path(entry).name}"
                out_path.write_bytes(z.read(entry))
                saved.append(out_path)
    except Exception as e:
        logger.warning("Erreur extraction images DOCX: %s", e)
    logger.info("%d image(s) extraite(s) depuis DOCX %s", len(saved), docx_path.name)
    return saved


def extract_fitz_text_pages(pdf_path: Path) -> str:
    """Extrait le texte brut page par page via PyMuPDF pour comparaison avec Docling."""
    sections: List[str] = []
    sections.append("## Extraction texte brut PyMuPDF/fitz — comparaison")
    sections.append("")
    sections.append("> Texte natif extrait directement par PyMuPDF sans traitement ML.")
    sections.append("")
    try:
        with fitz.open(str(pdf_path)) as doc:
            for page_num, page in enumerate(doc, start=1):
                text = page.get_text("text").strip()
                if not text:
                    continue
                sections.append(f"### Page {page_num}")
                sections.append("")
                sections.append("```text")
                sections.append(text)
                sections.append("```")
                sections.append("")
    except Exception as e:
        sections.append(f"Erreur PyMuPDF: {e}")
    return "\n".join(sections)


def fitz_extract_to_markdown(src_path: Path) -> str:
    """Extrait le texte PDF via PyMuPDF en mode 'words' pour reconstruire les espaces manquants.
    PyMuPDF analyse les positions des glyphes et insere les espaces a partir des gaps geometriques,
    ce qui gere mieux les PDFs ou les espaces ne sont pas encodes dans le flux texte."""
    if not src_path.is_file() or src_path.suffix.lower() != ".pdf":
        return f"fitz_direct: format non supporte ({src_path.suffix})"
    parts: List[str] = []
    try:
        with fitz.open(str(src_path)) as doc:
            for page_num, page in enumerate(doc, 1):
                parts.append(f"## Page {page_num}")
                parts.append("")
                words = page.get_text("words")  # [(x0,y0,x1,y1,word,block,line,word_no)]
                if not words:
                    parts.append("_(page sans texte)_")
                    parts.append("")
                    continue
                words.sort(key=lambda w: (w[1], w[0]))
                TOL_Y = 4.0
                current_y = words[0][1]
                current_block = words[0][5]
                line_buf: List[str] = []
                for w in words:
                    wy, wb = w[1], w[5]
                    new_line = abs(wy - current_y) > TOL_Y
                    new_block = wb != current_block
                    if new_line or new_block:
                        if line_buf:
                            parts.append(" ".join(line_buf))
                        if new_block:
                            parts.append("")
                        line_buf = [w[4]]
                        current_y = wy
                        current_block = wb
                    else:
                        line_buf.append(w[4])
                if line_buf:
                    parts.append(" ".join(line_buf))
                parts.append("")
    except Exception as e:
        parts.append(f"Erreur fitz: {e}")
    return "\n".join(parts)


def build_image_gallery_section(
    artifacts: List[ParsedArtifact],
    output_path: Path,
    max_images_per_pdf: int = 8,
) -> str:
    lines: List[str] = []
    gallery_root = output_path.parent / f"{output_path.stem}_assets"
    seen_sources = set()

    for artifact in artifacts:
        if artifact.source in seen_sources:
            continue
        seen_sources.add(artifact.source)

        if artifact.source_type != "file":
            continue

        src_path = Path(artifact.source)
        if src_path.suffix.lower() != ".pdf" or not src_path.is_file():
            continue

        pdf_gallery_dir = gallery_root / src_path.stem
        target_images = max(1, min(max_images_per_pdf, artifact.stats.get("images_refs", 12)))
        extracted_images = extract_figure_crops_to_dir(
            src_path,
            pdf_gallery_dir,
            max_images=target_images,
        )
        if not extracted_images:
            extracted_images = extract_pdf_images_to_dir(
                src_path,
                pdf_gallery_dir,
                max_images=target_images,
            )
        if not extracted_images:
            continue

        lines.append(f"## Galerie Images PDF: {src_path.name}")
        lines.append("")
        lines.append(
            f"Images extraites localement (fallback) pour affichage markdown: {len(extracted_images)}"
        )
        lines.append("")

        for idx, img_path in enumerate(extracted_images[:max_images_per_pdf], start=1):
            rel_path = img_path.relative_to(output_path.parent).as_posix()
            lines.append(f"### Image PDF {idx}")
            lines.append("")
            lines.append(f"![{img_path.name}]({rel_path})")
            lines.append("")

    return "\n".join(lines)


def gather_files_from_directory(folder: Path) -> List[Path]:
    files: List[Path] = []
    for p in folder.rglob("*"):
        if p.is_file():
            files.append(p)
    return sorted(files)


def likely_text_file(path: Path) -> bool:
    text_ext = {
        ".txt",
        ".md",
        ".markdown",
        ".rst",
        ".json",
        ".jsonl",
        ".yaml",
        ".yml",
        ".xml",
        ".html",
        ".htm",
        ".csv",
        ".tsv",
        ".py",
        ".js",
        ".ts",
        ".java",
        ".go",
        ".rs",
        ".cpp",
        ".c",
        ".h",
        ".hpp",
        ".sh",
        ".sql",
        ".tex",
        ".log",
    }
    if path.suffix.lower() in text_ext:
        return True
    guessed, _ = mimetypes.guess_type(path.name)
    return bool(guessed and guessed.startswith("text/"))


def read_text_fallback(path: Path) -> Optional[str]:
    if not likely_text_file(path):
        return None
    encodings = ["utf-8", "latin-1", "cp1252"]
    for enc in encodings:
        try:
            return path.read_text(encoding=enc, errors="strict")
        except Exception:
            continue
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return None


def build_docling_vlm_converter() -> DocumentConverter:
    if not VLM_AVAILABLE:
        raise RuntimeError("Pipeline VLM indisponible dans votre version Docling.")

    if not all([PdfFormatOption, InputFormat, PdfPipelineOptions]):
        raise RuntimeError("API Docling incomplete pour configurer le pipeline VLM.")

    if not OFFLINE_MODELS_DIR.is_dir():
        raise RuntimeError(f"Dossier de modeles local introuvable: {OFFLINE_MODELS_DIR}")

    if HAS_VLM_PRESET:
        vlm_options = VlmConvertOptions.from_preset("granite_docling")
        pipeline_opts = VlmPipelineOptions(
            vlm_options=vlm_options,
            artifacts_path=OFFLINE_MODELS_DIR,
            enable_remote_services=False,
            allow_external_plugins=False,
            generate_picture_images=False,
        )
    else:
        pipeline_opts = VlmPipelineOptions(
            vlm_options=vlm_model_specs.GRANITEDOCLING,
            artifacts_path=OFFLINE_MODELS_DIR,
            enable_remote_services=False,
            allow_external_plugins=False,
            generate_picture_images=False,
        )

    return DocumentConverter(
        format_options={
            InputFormat.PDF: PdfFormatOption(
                pipeline_cls=VlmPipeline,
                pipeline_options=pipeline_opts,
            )
        }
    )


def build_docling_standard_converter() -> DocumentConverter:
    if not all([PdfFormatOption, InputFormat, PdfPipelineOptions]):
        raise RuntimeError("API Docling incomplete pour configurer le pipeline standard.")

    if not OFFLINE_MODELS_DIR.is_dir():
        raise RuntimeError(f"Dossier de modeles local introuvable: {OFFLINE_MODELS_DIR}")

    pipeline_opts = PdfPipelineOptions()
    pipeline_opts.artifacts_path = OFFLINE_MODELS_DIR
    pipeline_opts.do_ocr = False
    pipeline_opts.do_table_structure = True

    # Mode TableFormer ACCURATE pour une meilleure qualite de parsing des tableaux
    try:
        from docling.datamodel.pipeline_options import TableFormerMode
        pipeline_opts.table_structure_options.mode = TableFormerMode.ACCURATE
        logger.info("Mode TableFormer ACCURATE active")
    except Exception:
        pass

    # pypdfium2 gere mieux les encodages de polices non-standard que docling-parse
    backend_cls = None
    try:
        from docling.backend.pypdfium2_backend import PyPdfiumDocumentBackend
        backend_cls = PyPdfiumDocumentBackend
        logger.info("Backend pypdfium2 actif (meilleur support encodages polices)")
    except Exception:
        logger.debug("pypdfium2 backend indisponible, utilisation du backend par defaut")

    try:
        fmt_option = (
            PdfFormatOption(pipeline_options=pipeline_opts, backend=backend_cls)
            if backend_cls is not None
            else PdfFormatOption(pipeline_options=pipeline_opts)
        )
    except Exception:
        fmt_option = PdfFormatOption(pipeline_options=pipeline_opts)

    return DocumentConverter(format_options={InputFormat.PDF: fmt_option})


def build_docling_standard_advanced_converter() -> DocumentConverter:
    if not all([PdfFormatOption, InputFormat, PdfPipelineOptions]):
        raise RuntimeError("API Docling incomplete pour configurer le pipeline standard avance.")

    if not OFFLINE_MODELS_DIR.is_dir():
        raise RuntimeError(f"Dossier de modeles local introuvable: {OFFLINE_MODELS_DIR}")

    pipeline_opts = PdfPipelineOptions()
    pipeline_opts.artifacts_path = OFFLINE_MODELS_DIR
    pipeline_opts.do_ocr = False
    pipeline_opts.do_table_structure = True
    # Profil plus pousse: enrichissement formules/code et pages image utiles au parsing complexe.
    pipeline_opts.do_formula_enrichment = True
    pipeline_opts.do_code_enrichment = True
    pipeline_opts.generate_page_images = True
    pipeline_opts.generate_picture_images = True

    return DocumentConverter(
        format_options={
            InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_opts),
        }
    )


def build_docling_standard_fast_converter() -> DocumentConverter:
    if not all([PdfFormatOption, InputFormat, PdfPipelineOptions]):
        raise RuntimeError("API Docling incomplete pour configurer le pipeline standard fast.")

    if not OFFLINE_MODELS_DIR.is_dir():
        raise RuntimeError(f"Dossier de modeles local introuvable: {OFFLINE_MODELS_DIR}")

    pipeline_opts = PdfPipelineOptions()
    pipeline_opts.artifacts_path = OFFLINE_MODELS_DIR
    pipeline_opts.enable_remote_services = False
    pipeline_opts.allow_external_plugins = False
    pipeline_opts.do_ocr = False
    pipeline_opts.do_table_structure = False
    pipeline_opts.do_picture_classification = False
    pipeline_opts.do_picture_description = False
    pipeline_opts.generate_page_images = False
    pipeline_opts.generate_picture_images = False
    pipeline_opts.force_backend_text = True

    return DocumentConverter(
        format_options={
            InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_opts),
        }
    )


def export_markdown_from_docling_document(doc: object) -> str:
    try:
        from docling.datamodel.document import ImageRefMode
        return doc.export_to_markdown(image_mode=ImageRefMode.PLACEHOLDER)
    except Exception:
        return doc.export_to_markdown()


def decode_html_entities(text: str) -> str:
    import html
    return html.unescape(text)


# Vocabulaire de reference pour detecter le decalage ASCII+29 (mots attendus apres decodage)
_GARBLED_VOCAB = frozenset([
    # Mots courants anglais
    "the", "and", "for", "are", "not", "with", "from", "this", "that", "have",
    "been", "will", "they", "more", "than", "each", "such", "only", "used",
    "shall", "must", "when", "where", "which", "also", "both", "then", "into",
    "upon", "after", "before", "during", "under", "over", "between", "within",
    "above", "below", "other", "these", "those", "their", "there", "here",
    "does", "made", "make", "take", "give", "work", "case", "time", "year",
    "part", "item", "step", "note", "date", "code", "task", "size", "data",
    "test", "open", "close", "next", "back", "left", "right", "side", "area",
    # Termes techniques generaux
    "general", "rules", "usage", "parts", "aircraft", "maintenance", "manual",
    "installation", "removal", "section", "figure", "table", "intended",
    "identify", "needed", "equipment", "information", "document", "description",
    "number", "revision", "issue", "page", "list", "type", "name", "unit",
    "purpose", "scope", "apply", "check", "record", "report", "design",
    "standard", "structure", "system", "perform", "following", "accordance",
    "drawing", "assembly", "repair", "inspection", "torque", "limit", "value",
    "sheet", "class", "level", "reference", "related", "required", "specified",
    "paragraph", "content", "status", "title", "weight", "engineering",
    "certification", "definitions", "preparation", "procedure", "attachment",
    "fastener", "sealant", "tools", "special", "access", "panel",
    "verify", "ensure", "confirm", "complete", "prior", "chapter",
    "function", "operation", "material", "method", "modification", "original",
    "production", "quality", "replacement", "requirement", "responsibility",
    "safety", "spare", "supplier", "support", "technical", "validation",
    "verification", "version", "applicable", "application", "authorizer",
    "authorization", "classification", "component", "condition", "impact",
    "include", "install", "instruction", "compliance", "control", "criteria",
    "concession", "assessment", "assembly", "applicable", "category",
    "effectivity", "tolerance", "specification", "configuration", "compliance",
    "airworthiness", "release", "approval", "approved", "issue", "status",
    "chapter", "paragraph", "subparagraph", "subject", "title", "contents",
    "introduction", "background", "scope", "applicability", "definitions",
    "references", "requirements", "procedures", "documentation", "records",
    "reporting", "training", "responsibilities", "management", "planning",
    "implementation", "monitoring", "review", "actions", "corrective",
    "preventive", "audit", "findings", "evidence", "objective", "result",
])


def _detect_garbled_text(md: str, shift: int = 29) -> bool:
    """Detecte si le texte extrait contient un decalage ASCII systematique."""
    tokens = re.findall(r'[!"#$%&\'()*+,\-./0-9:;<=>?@A-Z\[\\\]^_]{4,}', md[:5000])
    if len(tokens) < 10:
        return False
    hits = 0
    for tok in tokens[:60]:
        decoded = ''.join(
            chr(ord(c) + shift) if ord(c) + shift <= 126 else c for c in tok
        ).lower()
        if any(word in decoded for word in _GARBLED_VOCAB if len(word) >= 4):
            hits += 1
    ratio = hits / min(len(tokens), 60)
    if ratio > 0.12:
        logger.info("Encodage decale detecte (ratio=%.2f), correction +%d appliquee", ratio, shift)
        return True
    return False


def _shift_str(s: str, shift: int = 29) -> str:
    """Applique le decalage +29 uniquement aux caracteres dans la plage garbled (33-97).
    Preserves les espaces, tabulations, accents et caracteres > 97."""
    result = []
    for c in s:
        code = ord(c)
        if 33 <= code <= 97:
            new_code = code + shift
            if new_code <= 126:
                result.append(chr(new_code))
                continue
        result.append(c)
    return "".join(result)


def _content_is_garbled(text: str, shift: int = 29) -> bool:
    """Detecte si un fragment de texte est garble (decalage +29).
    Contrairement a _detect_garbled_text, fonctionne sur un seul fragment court."""
    tokens = re.findall(r'[!"#$%&\'()*+,\-./0-9:;<=>?@A-Z\[\\\]^_]{4,}', text)
    if not tokens:
        return False
    hits = sum(
        1 for tok in tokens
        if any(
            w in "".join(chr(ord(c) + shift) if ord(c) + shift <= 126 else c for c in tok).lower()
            for w in _GARBLED_VOCAB if len(w) >= 4
        )
    )
    return hits > 0 and hits / len(tokens) >= 0.15


def fix_garbled_text_if_needed(md: str) -> str:
    """Corrige le decalage ASCII+29 ligne par ligne.
    Seules les lignes detectees comme garblees sont modifiees — le texte deja correct est preserve."""
    if not FIX_GARBLED_TEXT or not _detect_garbled_text(md):
        return md

    table_sep_re = re.compile(r"^\s*\|[-|: ]+\|\s*$")
    hr_re = re.compile(r"^\s*[-*_]{3,}\s*$")
    lines_out: List[str] = []
    in_code_block = False

    for line in md.split("\n"):
        stripped = line.strip()

        if stripped.startswith("```") or stripped.startswith("~~~"):
            in_code_block = not in_code_block
            lines_out.append(line)
            continue
        if in_code_block or not stripped:
            lines_out.append(line)
            continue
        if stripped.startswith("<"):
            lines_out.append(line)
            continue
        if hr_re.match(stripped) or table_sep_re.match(stripped):
            lines_out.append(line)
            continue

        # Titre : corriger le contenu uniquement si garble
        m = re.match(r"^(#{1,6} )(.*)", line)
        if m:
            content = m.group(2)
            lines_out.append(m.group(1) + (_shift_str(content) if _content_is_garbled(content) else content))
            continue

        # Ligne de tableau : corriger cellule par cellule
        if stripped.startswith("|") and stripped.endswith("|"):
            parts = line.split("|")
            lines_out.append("|".join(
                _shift_str(p) if 0 < i < len(parts) - 1 and _content_is_garbled(p) else p
                for i, p in enumerate(parts)
            ))
            continue

        # Liste
        m = re.match(r"^(\s*(?:[-*+]|\d+\.)\s+)(.*)", line)
        if m:
            content = m.group(2)
            lines_out.append(m.group(1) + (_shift_str(content) if _content_is_garbled(content) else content))
            continue

        # Blockquote
        m = re.match(r"^(>+\s*)(.*)", line)
        if m:
            content = m.group(2)
            lines_out.append(m.group(1) + (_shift_str(content) if _content_is_garbled(content) else content))
            continue

        # Texte ordinaire
        lines_out.append(_shift_str(line) if _content_is_garbled(line) else line)

    return "\n".join(lines_out)


def convert_doc_to_docx(doc_path: Path) -> Optional[Path]:
    """Convertit un fichier .doc binaire en .docx via LibreOffice headless."""
    docx_path = doc_path.with_suffix(".docx")
    if docx_path.exists():
        logger.info("Fichier .docx deja present: %s", docx_path.name)
        return docx_path
    import subprocess
    candidates = [
        "soffice",
        "libreoffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for exe in candidates:
        try:
            subprocess.run(
                [exe, "--headless", "--convert-to", "docx",
                 "--outdir", str(doc_path.parent), str(doc_path)],
                capture_output=True,
                timeout=120,
            )
            if docx_path.exists():
                logger.info("Conversion .doc -> .docx reussie: %s", docx_path.name)
                return docx_path
        except (FileNotFoundError, OSError):
            continue
        except Exception:
            continue
    logger.warning("Conversion .doc -> .docx echouee pour %s (LibreOffice requis)", doc_path.name)
    return None


def parse_single_source(
    source: str,
    method: str,
    include_image_contexts: bool,
    standard_fast_mode: bool,
    converter: Optional[DocumentConverter] = None,
) -> ParsedArtifact:
    source_type = detect_source_type(source)
    start = time.time()
    warnings: List[str] = []

    # Auto-conversion .doc -> .docx (Docling ne supporte pas le format .doc binaire)
    if source_type == "file" and Path(source).suffix.lower() == ".doc":
        converted = convert_doc_to_docx(Path(source))
        if converted is not None:
            source = str(converted)
            warnings.append(f"Fichier .doc converti automatiquement en .docx: {converted.name}")
        else:
            warnings.append("Conversion .doc -> .docx impossible (LibreOffice non trouve)")
            # On continue quand meme; Docling va echouer et le fallback texte sera tente

    if source_type not in {"file", "url"}:
        return ParsedArtifact(
            source=source,
            source_type=source_type,
            method="unsupported",
            parse_elapsed_s=0,
            markdown="",
            stats=analyze_markdown_content(""),
            image_contexts=[],
            warnings=["Source non supportee"],
            error="Seuls les fichiers et URLs sont pris en charge ici.",
        )

    if method not in SUPPORTED_METHODS:
        return ParsedArtifact(
            source=source,
            source_type=source_type,
            method=method,
            parse_elapsed_s=0,
            markdown="",
            stats=analyze_markdown_content(""),
            image_contexts=[],
            warnings=["Methode non supportee"],
            error=f"Methode invalide: {method}",
        )

    # Methode fitz_direct: bypass Docling, extraction native PyMuPDF
    if method == "fitz_direct":
        src = Path(source)
        md = fitz_extract_to_markdown(src)
        elapsed = time.time() - start
        return ParsedArtifact(
            source=source,
            source_type=source_type,
            method=method,
            parse_elapsed_s=elapsed,
            markdown=md,
            stats=analyze_markdown_content(md),
            image_contexts=[],
            warnings=warnings,
        )

    if converter is None:
        if method == "docling_standard":
            converter = (
                build_docling_standard_fast_converter()
                if standard_fast_mode
                else build_docling_standard_converter()
            )
        elif method == "docling_standard_advanced":
            converter = build_docling_standard_advanced_converter()
        else:
            converter = build_docling_vlm_converter()

    try:
        result = converter.convert(source)
        md = fix_garbled_text_if_needed(
            decode_html_entities(export_markdown_from_docling_document(result.document))
        )
        elapsed = time.time() - start
        stats = analyze_markdown_content(md)
        contexts = (
            extract_image_contexts(md, context_window=DEFAULT_CONTEXT_WINDOW)
            if include_image_contexts
            else []
        )
        free_memory()
        return ParsedArtifact(
            source=source,
            source_type=source_type,
            method=method,
            parse_elapsed_s=elapsed,
            markdown=md,
            stats=stats,
            image_contexts=contexts,
            warnings=warnings,
        )
    except Exception as docling_error:
        if source_type == "file":
            fallback_text = read_text_fallback(Path(source))
            if fallback_text is not None:
                warnings.append(
                    "Echec Docling sur ce fichier; fallback lecture texte brute active."
                )
                elapsed = time.time() - start
                stats = analyze_markdown_content(fallback_text)
                contexts = (
                    extract_image_contexts(fallback_text, context_window=DEFAULT_CONTEXT_WINDOW)
                    if include_image_contexts
                    else []
                )
                free_memory()
                return ParsedArtifact(
                    source=source,
                    source_type=source_type,
                    method=f"{method}+text_fallback",
                    parse_elapsed_s=elapsed,
                    markdown=fallback_text,
                    stats=stats,
                    image_contexts=contexts,
                    warnings=warnings,
                )

        elapsed = time.time() - start
        free_memory()
        return ParsedArtifact(
            source=source,
            source_type=source_type,
            method=method,
            parse_elapsed_s=elapsed,
            markdown="",
            stats=analyze_markdown_content(""),
            image_contexts=[],
            warnings=warnings,
            error=str(docling_error),
        )


def parse_anything(
    source: str,
    methods: List[str],
    include_image_contexts: bool,
    standard_fast_mode: bool,
) -> List[ParsedArtifact]:
    src_type = detect_source_type(source)
    artifacts: List[ParsedArtifact] = []

    if src_type == "directory":
        files = gather_files_from_directory(Path(source))
        logger.info("Dossier detecte: %s fichiers a traiter", len(files))
        for file_path in files:
            logger.info("Traitement: %s", file_path)
            converters: Dict[str, Optional[DocumentConverter]] = {}
            for method in methods:
                if method == "fitz_direct":
                    pass  # pas de converter Docling
                elif method == "docling_standard":
                    if method not in converters:
                        try:
                            converters[method] = build_docling_standard_converter()
                        except Exception as e:
                            logger.warning("Init converter standard impossible: %s", e)
                            converters[method] = None
                elif method == "docling_standard_advanced":
                    if method not in converters:
                        try:
                            converters[method] = build_docling_standard_advanced_converter()
                        except Exception as e:
                            logger.warning("Init converter standard avance impossible: %s", e)
                            converters[method] = None
                elif method == "docling_vlm_granite":
                    if method not in converters:
                        try:
                            converters[method] = build_docling_vlm_converter()
                        except Exception as e:
                            logger.warning("Init converter VLM impossible: %s", e)
                            converters[method] = None

                artifact = parse_single_source(
                    str(file_path),
                    method=method,
                    include_image_contexts=include_image_contexts,
                    standard_fast_mode=standard_fast_mode,
                    converter=converters.get(method),
                )
                artifacts.append(artifact)
    else:
        for method in methods:
            artifact = parse_single_source(
                source,
                method=method,
                include_image_contexts=include_image_contexts,
                standard_fast_mode=standard_fast_mode,
            )
            artifacts.append(artifact)

    return artifacts


def benchmark_methods(
    source: str,
    methods: List[str],
    runs_per_method: int,
    include_image_contexts: bool,
) -> List[MethodBenchmark]:
    benches: List[MethodBenchmark] = []
    for method in methods:
        label = METHOD_LABELS.get(method, method)
        logger.info("=== Benchmark %s (%s runs) ===", label, runs_per_method)
        runs: List[ParsedArtifact] = []

        converter: Optional[DocumentConverter] = None
        if method == "fitz_direct":
            pass  # pas de converter Docling
        elif method == "docling_standard":
            converter = build_docling_standard_converter()
        elif method == "docling_standard_advanced":
            converter = build_docling_standard_advanced_converter()
        elif method == "docling_vlm_granite":
            converter = build_docling_vlm_converter()

        for run_idx in range(1, runs_per_method + 1):
            logger.info("Run %s/%s pour %s", run_idx, runs_per_method, label)
            artifact = parse_single_source(
                source=source,
                method=method,
                include_image_contexts=include_image_contexts,
                standard_fast_mode=False,
                converter=converter,
            )
            artifact = replace(
                artifact,
                warnings=artifact.warnings + [f"benchmark_run:{run_idx}/{runs_per_method}"],
            )
            runs.append(artifact)

        benches.append(MethodBenchmark(method=method, label=label, runs=runs))

    return benches


def render_markdown_report(
    source: str,
    benchmarks: List[MethodBenchmark],
    include_image_contexts: bool,
    include_raw_markdown: bool,
    context_window: int = DEFAULT_CONTEXT_WINDOW,
) -> str:
    artifacts = [b.best for b in benchmarks]
    success = [a for a in artifacts if not a.error]
    failures = [a for a in artifacts if a.error]

    lines: List[str] = []
    lines.append("# Rapport RAG multi-formats (Docling)")
    lines.append("")
    lines.append(f"Source d'entree: {source}")
    lines.append(f"Date: {time.strftime('%Y-%m-%d %H:%M:%S')}")
    lines.append("")
    lines.append("## Parametres")
    lines.append("")
    lines.append("- Execution en dur: aucune option CLI necessaire")
    lines.append(f"- Fichier source fixe: {source}")
    lines.append(f"- Repetitions par methode: {RUNS_PER_METHOD}")
    lines.append("- OCR: desactive explicitement pour les images")
    lines.append(f"- Fenetre de contexte image: +/- {context_window} caracteres")
    lines.append("- Envoi LLM: desactive (simulation seulement)")
    lines.append("")

    lines.append("## Resume global")
    lines.append("")
    lines.append(f"- Total elements traites: {len(artifacts)}")
    lines.append(f"- Succes: {len(success)}")
    lines.append(f"- Echecs: {len(failures)}")
    lines.append(f"- Images detectees: {sum(a.stats.get('images_refs', 0) for a in artifacts)}")
    lines.append("")

    lines.append("## Benchmark Temps (multi-runs)")
    lines.append("")
    lines.append("| Methode | Run 1 (s) | Run 2 (s) | Moyenne (s) | Meilleur parse (s) |")
    lines.append("|---|---:|---:|---:|---:|")
    for b in benchmarks:
        run_values = [f"{r.parse_elapsed_s:.2f}" for r in b.runs]
        run1 = run_values[0] if len(run_values) > 0 else "-"
        run2 = run_values[1] if len(run_values) > 1 else "-"
        best_s = min(r.parse_elapsed_s for r in b.runs) if b.runs else 0.0
        lines.append(
            f"| {b.label} | {run1} | {run2} | {b.avg_parse_s:.2f} | {best_s:.2f} |"
        )
    lines.append("")

    lines.append("## Tableau des resultats")
    lines.append("")
    lines.append("| Source | Methode | Temps Docling parse (s) | Mots | Tables | Formules bloc | Images | Statut |")
    lines.append("|---|---|---:|---:|---:|---:|---:|---|")
    for a in artifacts:
        status = "OK" if not a.error else "ERREUR"
        lines.append(
            "| "
            f"{a.source} | {a.method} | {a.parse_elapsed_s:.2f} | {a.stats.get('words', 0)} | "
            f"{a.stats.get('tables_count', 0)} | {a.stats.get('math_block', 0)} | "
            f"{a.stats.get('images_refs', 0)} | {status} |"
        )
    lines.append("")

    for idx, a in enumerate(artifacts, start=1):
        lines.append(f"## Element {idx}: {a.source}")
        lines.append("")
        lines.append(f"- Type: {a.source_type}")
        lines.append(f"- Methode: {a.method}")
        lines.append(f"- Temps Docling parse pur: {a.parse_elapsed_s:.2f}s")
        if a.error:
            lines.append(f"- Erreur: {a.error}")
        if a.warnings:
            lines.append("- Warnings:")
            for w in a.warnings:
                lines.append(f"  - {w}")
        lines.append("")

        lines.append("### Metriques")
        lines.append("")
        lines.append("```json")
        lines.append(json.dumps(a.stats, ensure_ascii=False, indent=2))
        lines.append("```")
        lines.append("")

        lines.append("### Images")
        lines.append("")
        lines.append(f"Placeholders images Docling detectes: {a.stats.get('images_refs', 0)}")
        lines.append("(Les images reelles sont extraites et affichees dans le fichier `_extracted.md`)")
        lines.append("")

        if include_raw_markdown:
            lines.append("### Contenu extrait (rendu)")
            lines.append("")
            lines.append("---")
            lines.append("")
            lines.append(a.markdown)
            lines.append("")
            lines.append("---")
            lines.append("")

    return "\n".join(lines)


def inject_images_into_markdown(
    md: str,
    src_path: Path,
    md_output_path: Path,
    max_images: int = 30,
) -> str:
    """Extrait les images du document source et les insere a la fin du markdown.
    Supporte PDF (PyMuPDF), PPTX (python-pptx), DOCX (zipfile)."""
    assets_dir = md_output_path.parent / f"{md_output_path.stem}_assets"
    suffix = src_path.suffix.lower()
    images: List[Path] = []

    if suffix == ".pdf":
        embedded = extract_pdf_images_to_dir(src_path, assets_dir, max_images=max_images)
        crops = extract_figure_crops_to_dir(src_path, assets_dir, max_images=max_images)
        seen: set = set()
        for img in embedded + crops:
            if img.name not in seen:
                seen.add(img.name)
                images.append(img)
        logger.info("%d image(s) PDF (%d raster + %d crops) depuis %s",
                    len(images), len(embedded), len(crops), src_path.name)
    elif suffix in {".pptx", ".ppt"}:
        images = extract_pptx_images_to_dir(src_path, assets_dir, max_images=max_images)
    elif suffix in {".docx", ".doc"}:
        images = extract_docx_images_to_dir(src_path, assets_dir, max_images=max_images)

    images = images[:max_images]

    if not images:
        logger.warning("Aucune image extraite depuis: %s", src_path.name)
        return md

    lines = [md, "", "---", "", f"## Images extraites ({src_path.name})", ""]
    for i, img_path in enumerate(images, 1):
        rel = img_path.relative_to(md_output_path.parent).as_posix()
        lines.append(f"### Figure {i}")
        lines.append("")
        lines.append(f"![Figure {i}: {img_path.name}]({rel})")
        lines.append("")
    return "\n".join(lines)


def main() -> None:
    source = INPUT_SOURCE
    output = OUTPUT_MD
    methods = ["docling_standard", "fitz_direct"]

    logger.info("Demarrage pipeline RAG")
    logger.info("Source: %s", source)
    logger.info("OCR force pages completes: %s", FORCE_FULL_PAGE_OCR)

    benchmarks = benchmark_methods(
        source=source,
        methods=methods,
        runs_per_method=RUNS_PER_METHOD,
        include_image_contexts=INCLUDE_IMAGE_CONTEXTS,
    )

    report = render_markdown_report(
        source=source,
        benchmarks=benchmarks,
        include_image_contexts=INCLUDE_IMAGE_CONTEXTS,
        include_raw_markdown=INCLUDE_RAW_MARKDOWN,
        context_window=DEFAULT_CONTEXT_WINDOW,
    )

    if INCLUDE_PDF_GALLERY:
        artifacts_for_gallery = [b.best for b in benchmarks]
        gallery_section = build_image_gallery_section(
            artifacts_for_gallery,
            output,
            max_images_per_pdf=max(1, int(GALLERY_MAX_IMAGES)),
        )
        if gallery_section:
            report = report + "\n\n" + gallery_section + "\n"

    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(report, encoding="utf-8")
    logger.info("Rapport genere: %s", output)

    src_path = Path(source)
    for b in benchmarks:
        best = b.best
        if not best.markdown or best.error:
            continue
        src_stem = src_path.stem if best.source_type == "file" else "source"
        clean_name = f"{src_stem}_{best.method.replace('+', '_')}_extracted.md"
        clean_path = output.parent / clean_name
        md = best.markdown

        # Injection d'images uniquement pour la methode docling (pas fitz_direct)
        if best.method != "fitz_direct" and best.source_type == "file" and src_path.suffix.lower() in {".pdf", ".pptx", ".ppt", ".docx", ".doc"}:
            md = inject_images_into_markdown(md, src_path, clean_path, max_images=IMAGES_PER_PDF)

        clean_path.write_text(md, encoding="utf-8")
        logger.info("Markdown extrait sauvegarde: %s", clean_path)


if __name__ == "__main__":
    main()
