import io
import json
import os
import re
import shutil
import subprocess
import tempfile
import time
from importlib import metadata as importlib_metadata
from pathlib import Path
from typing import Any, Callable, Dict, Iterable, List, Optional

import pandas as pd
import pdfplumber
from docx import Document
from docx.oxml.ns import qn
from lxml import etree
from pptx import Presentation

try:
    import numpy as np
except Exception:  # pragma: no cover - optional dependency
    np = None

try:
    from PIL import Image, ImageOps
except Exception:  # pragma: no cover - optional dependency
    Image = None
    ImageOps = None

try:
    import pypdfium2 as pdfium
except Exception:  # pragma: no cover - optional dependency
    pdfium = None

try:
    import pytesseract
except Exception:  # pragma: no cover - optional dependency
    pytesseract = None

try:
    from rapidocr_onnxruntime import RapidOCR
except Exception:  # pragma: no cover - optional dependency
    RapidOCR = None

from .constants import (
    ANTIWORD_BIN,
    ANTIWORD_CPU_SECONDS,
    ANTIWORD_MEMORY_MB,
    ANTIWORD_SHARE_DIR,
    EXCEL_INCLUDE_ROW_CONTEXT,
    EXCEL_TABLE_PREVIEW_ROWS,
    LOCAL_PARSE_MAX_RETRIES,
    OCR_ENABLED,
    OCR_ENGINE_PRIORITY,
    OCR_IMAGE_PREPROCESS,
    OCR_MAX_IMAGE_SIDE_PX,
    OCR_MAX_PDF_PAGES,
    OCR_MIN_CONFIDENCE,
    OCR_MIN_IMAGE_SIDE_PX,
    OCR_MIN_PDF_PAGE_TEXT_CHARS,
    OCR_PDF_FALLBACK_ENABLED,
    OCR_PDF_RENDER_DPI,
    OCR_TESSERACT_BIN,
    OCR_TESSERACT_LANG,
    OCR_TESSERACT_OEM,
    OCR_TESSERACT_PSM,
    OCR_TIMEOUT_SECONDS,
    PDF_EXCLUDE_TABLE_TEXT,
    MAX_CHUNK_CHARS,
    MAX_CSV_ROWS,
    MAX_DOCUMENT_BYTES,
    MAX_ERROR_MESSAGE_CHARS,
    MAX_EXCEL_ROWS_PER_SHEET,
    MAX_EXCEL_SHEETS,
    MAX_FALLBACK_TEXT_CHARS,
    MAX_PDF_PAGES,
    MAX_PPTX_SLIDES,
    PARSE_TIMEOUT_SECONDS,
    TEXTLIKE_EXTENSIONS,
)
from .retry import is_retryable_error
from .text_utils import normalize_text, rows_to_markdown_table, safe_decode, strip_html_tags

_LIB_VERSION_CACHE: Dict[str, str] = {}
_BINARY_SIGNATURE_FORMATS = {"pdf", "jpeg", "png", "tiff", "webp", "ole2", "ooxml", "rtf"}
_IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "tif", "tiff", "bmp", "webp"}
_RAPID_OCR_ENGINE: Optional[Any] = None


def _response(text: str, parser_error: Optional[str], parser_strategy: str) -> Dict[str, Any]:
    return {
        "text": text,
        "parser_error": parser_error,
        "parser_strategy": parser_strategy,
        "ocr_attempted": False,
        "ocr_used": False,
        "ocr_engine_trace": None,
        "ocr_pages": 0,
        "ocr_supplement_pages": 0,
    }


def _with_ocr_metadata(
    result: Dict[str, Any],
    *,
    attempted: bool,
    used: bool,
    trace: Optional[str],
    pages: int,
    supplement_pages: int,
) -> Dict[str, Any]:
    result["ocr_attempted"] = bool(attempted)
    result["ocr_used"] = bool(used)
    result["ocr_engine_trace"] = normalize_text(trace or "") or None
    result["ocr_pages"] = max(0, _safe_int(pages))
    result["ocr_supplement_pages"] = max(0, _safe_int(supplement_pages))
    return result


def _sanitize_error_message(message: str) -> str:
    if not message:
        return ""

    cleaned = str(message).replace("\r", " ").replace("\n", " ").strip()
    cleaned = re.sub(r"[A-Za-z]:\\[^\s]+", "<path>", cleaned)
    cleaned = re.sub(r"/(?:[^/\s]+/)*[^/\s]+", "<path>", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned)

    if len(cleaned) > MAX_ERROR_MESSAGE_CHARS:
        cleaned = f"{cleaned[: MAX_ERROR_MESSAGE_CHARS - 3]}..."
    return cleaned


def _error(code: str, message: str) -> str:
    return f"{code}: {_sanitize_error_message(message)}"


def _is_timed_out(started_at: float, timeout_seconds: float = PARSE_TIMEOUT_SECONDS) -> bool:
    return (time.time() - started_at) > timeout_seconds


def _normalize_extension(extension: str) -> str:
    ext = (extension or "").lower().strip()
    return ext[1:] if ext.startswith(".") else ext


def _detect_binary_format(content: bytes) -> str:
    if not content or len(content) < 5:
        return "unknown"
    if content[:5] == b"%PDF-":
        return "pdf"
    if content[:3] == b"\xff\xd8\xff":
        return "jpeg"
    if content[:4] == b"\x89PNG":
        return "png"
    if content[:4] in {b"II*\x00", b"MM\x00*"}:
        return "tiff"
    if content[:4] == b"RIFF" and content[8:12] == b"WEBP":
        return "webp"
    if content[:4] == b"\xd0\xcf\x11\xe0":
        return "ole2"
    if content[:4] == b"PK\x03\x04":
        return "ooxml"
    if content.lstrip()[:5] == b"{\\rtf":
        return "rtf"
    return "unknown"


def _is_mime_mismatch_for_textlike(extension: str, detected_format: str) -> bool:
    ext = _normalize_extension(extension)
    if ext not in TEXTLIKE_EXTENSIONS:
        return False
    return detected_format in (_BINARY_SIGNATURE_FORMATS - {"rtf"}) or (detected_format == "rtf" and ext != "rtf")


def _truncate_text(text: str, max_chars: int = MAX_FALLBACK_TEXT_CHARS) -> str:
    if len(text) <= max_chars:
        return text
    return f"{text[:max_chars]}\n\n[TRUNCATED at {max_chars} chars]"


def _lib_version(lib_name: str) -> str:
    if lib_name in _LIB_VERSION_CACHE:
        return _LIB_VERSION_CACHE[lib_name]

    try:
        _LIB_VERSION_CACHE[lib_name] = importlib_metadata.version(lib_name)
    except Exception:
        _LIB_VERSION_CACHE[lib_name] = "unknown"
    return _LIB_VERSION_CACHE[lib_name]


def _with_lib_versions(strategy: str, *libs: str) -> str:
    if not libs:
        return strategy
    versions = [f"{lib}={_lib_version(lib)}" for lib in libs]
    return f"{strategy}|libs:{','.join(versions)}"


def _strip_rtf_control_words(decoded: str) -> str:
    # Keep plain textual payload while dropping most RTF control markup.
    no_hex = re.sub(r"\\'[0-9a-fA-F]{2}", " ", decoded)
    no_controls = re.sub(r"\\[a-zA-Z]+-?\d*\s?", " ", no_hex)
    no_groups = no_controls.replace("{", " ").replace("}", " ")
    return normalize_text(no_groups)


def _extract_docx_comments(document: Document) -> List[str]:
    comments: List[str] = []
    comments_part = getattr(document.part, "comments_part", None)
    if comments_part is None:
        return comments

    root = getattr(comments_part, "element", None)
    if root is None:
        return comments

    try:
        for comment in root.xpath(".//*[local-name()='comment']"):
            author = normalize_text(comment.get(qn("w:author"), ""))
            text_fragments = [normalize_text(node.text or "") for node in comment.xpath(".//*[local-name()='t']")]
            text = normalize_text(" ".join([fragment for fragment in text_fragments if fragment]))
            if text:
                label = f"[{author}] " if author else ""
                comments.append(f"{label}{text}")
    except Exception:
        return []

    return comments


def _extract_docx_textboxes(document: Document) -> List[str]:
    seen = set()
    extracted: List[str] = []
    for node in document.element.body.xpath(".//*[local-name()='txbxContent']//*[local-name()='t']"):
        text = normalize_text(node.text or "")
        if text and text not in seen:
            seen.add(text)
            extracted.append(text)
    return extracted


def _collect_docx_image_blobs(document: Document) -> List[bytes]:
    blobs: List[bytes] = []
    seen = set()

    def _collect_from_part(part: Any) -> None:
        rels = getattr(part, "rels", {})
        for rel in rels.values():
            reltype = str(getattr(rel, "reltype", ""))
            if "image" not in reltype:
                continue

            target_part = getattr(rel, "target_part", None)
            blob = getattr(target_part, "blob", None)
            if not blob:
                continue

            signature = (len(blob), blob[:32])
            if signature in seen:
                continue

            seen.add(signature)
            blobs.append(blob)

    _collect_from_part(document.part)
    for section in document.sections:
        try:
            _collect_from_part(section.header.part)
        except Exception:
            pass
        try:
            _collect_from_part(section.footer.part)
        except Exception:
            pass

    return blobs


def _extract_paragraph_with_hyperlinks(paragraph) -> str:
    rels = paragraph.part.rels
    pieces: List[str] = []

    for child in paragraph._p.iterchildren():
        if child.tag == qn("w:r"):
            run_text = "".join([node.text for node in child.xpath(".//*[local-name()='t']") if node.text])
            normalized = normalize_text(run_text)
            if normalized:
                pieces.append(normalized)
        elif child.tag == qn("w:hyperlink"):
            link_text = "".join([node.text for node in child.xpath(".//*[local-name()='t']") if node.text])
            rid = child.get(qn("r:id"))
            href = ""
            if rid and rid in rels:
                href = normalize_text(str(getattr(rels[rid], "target_ref", "")))

            normalized_link_text = normalize_text(link_text)
            if normalized_link_text and href:
                pieces.append(f"{normalized_link_text} ({href})")
            elif normalized_link_text:
                pieces.append(normalized_link_text)
            elif href:
                pieces.append(href)

    text = normalize_text(" ".join(pieces))
    return text or normalize_text(paragraph.text)


def _antiword_preexec_limits() -> Optional[Callable[[], None]]:
    if os.name != "posix":
        return None

    try:
        import resource
    except ImportError:
        return None

    def _apply_limits() -> None:
        try:
            resource.setrlimit(resource.RLIMIT_CPU, (ANTIWORD_CPU_SECONDS, ANTIWORD_CPU_SECONDS))
        except Exception:
            pass

        try:
            max_memory_bytes = ANTIWORD_MEMORY_MB * 1024 * 1024
            resource.setrlimit(resource.RLIMIT_AS, (max_memory_bytes, max_memory_bytes))
        except Exception:
            pass

    return _apply_limits


def _excel_engine_for_extension(extension: str) -> Optional[str]:
    ext = _normalize_extension(extension)
    if ext == "xls":
        return "xlrd"
    if ext in {"xlsx", "xlsm"}:
        return "openpyxl"
    if ext == "xlsb":
        return "pyxlsb"
    return None


def _safe_int(value: Any) -> int:
    try:
        return int(value)
    except Exception:
        return 0


def _point_inside_bbox(x: float, y: float, bbox: tuple) -> bool:
    if not bbox or len(bbox) != 4:
        return False
    x0, top, x1, bottom = bbox
    return float(x0) <= x <= float(x1) and float(top) <= y <= float(bottom)


def _words_to_lines(words: List[Dict[str, Any]], y_tolerance: float = 3.0) -> str:
    if not words:
        return ""

    sorted_words = sorted(
        words,
        key=lambda word: (float(word.get("top", 0.0)), float(word.get("x0", 0.0))),
    )
    lines: List[str] = []
    current_words: List[str] = []
    current_top: Optional[float] = None

    for word in sorted_words:
        text = normalize_text(str(word.get("text", "")))
        if not text:
            continue

        word_top = float(word.get("top", 0.0))
        if current_top is None or abs(word_top - current_top) <= y_tolerance:
            current_words.append(text)
            if current_top is None:
                current_top = word_top
            else:
                current_top = (current_top + word_top) / 2.0
            continue

        if current_words:
            lines.append(" ".join(current_words))
        current_words = [text]
        current_top = word_top

    if current_words:
        lines.append(" ".join(current_words))

    return normalize_text("\n".join(lines))


def _extract_pdf_text_excluding_tables(page, table_bboxes: List[tuple]) -> str:
    try:
        words = page.extract_words(
            x_tolerance=2,
            y_tolerance=2,
            keep_blank_chars=False,
            use_text_flow=True,
        )
    except Exception:
        words = []

    if not words:
        return ""

    filtered_words = []
    for word in words:
        x0 = float(word.get("x0", 0.0))
        x1 = float(word.get("x1", x0))
        top = float(word.get("top", 0.0))
        bottom = float(word.get("bottom", top))
        mid_x = (x0 + x1) / 2.0
        mid_y = (top + bottom) / 2.0

        if any(_point_inside_bbox(mid_x, mid_y, bbox) for bbox in table_bboxes):
            continue
        filtered_words.append(word)

    return _words_to_lines(filtered_words)


def _resize_image_for_ocr(image):
    width, height = image.size
    if width <= 0 or height <= 0:
        return image

    min_side = min(width, height)
    max_side = max(width, height)
    scale = 1.0

    if min_side < OCR_MIN_IMAGE_SIDE_PX:
        scale = OCR_MIN_IMAGE_SIDE_PX / float(min_side)
    if max_side * scale > OCR_MAX_IMAGE_SIDE_PX:
        scale = OCR_MAX_IMAGE_SIDE_PX / float(max_side)

    if abs(scale - 1.0) < 0.05:
        return image

    new_size = (
        max(1, int(round(width * scale))),
        max(1, int(round(height * scale))),
    )
    return image.resize(new_size)


def _prepare_image_for_ocr(raw_image):
    image = raw_image
    if ImageOps is not None:
        image = ImageOps.exif_transpose(image)

    if image.mode not in {"L", "RGB"}:
        image = image.convert("RGB")

    if OCR_IMAGE_PREPROCESS and ImageOps is not None:
        image = ImageOps.grayscale(image)
        image = ImageOps.autocontrast(image)

    image = _resize_image_for_ocr(image)
    return image


def _polygon_anchor(box: Any) -> tuple:
    if isinstance(box, (list, tuple)) and box:
        first_point = box[0]
        if isinstance(first_point, (list, tuple)) and len(first_point) >= 2:
            return float(first_point[1]), float(first_point[0])
    return 0.0, 0.0


def _parse_rapidocr_result(result: Any) -> tuple:
    lines_with_anchor = []
    scores = []

    if not isinstance(result, list):
        return "", 0.0, 0

    for item in result:
        if not isinstance(item, (list, tuple)) or len(item) < 3:
            continue

        text = normalize_text(str(item[1]))
        if not text:
            continue

        try:
            score = float(item[2])
        except Exception:
            score = 0.0

        if score < OCR_MIN_CONFIDENCE:
            continue

        anchor_y, anchor_x = _polygon_anchor(item[0])
        lines_with_anchor.append((anchor_y, anchor_x, text))
        scores.append(score)

    if not lines_with_anchor:
        return "", 0.0, 0

    lines_with_anchor.sort(key=lambda item: (item[0], item[1]))
    ordered_lines = [line for _, _, line in lines_with_anchor]
    avg_score = sum(scores) / len(scores) if scores else 0.0
    return normalize_text("\n".join(ordered_lines)), avg_score, len(ordered_lines)


def _get_rapidocr_engine() -> Optional[Any]:
    global _RAPID_OCR_ENGINE

    if RapidOCR is None:
        return None
    if _RAPID_OCR_ENGINE is not None:
        return _RAPID_OCR_ENGINE

    try:
        _RAPID_OCR_ENGINE = RapidOCR()
    except Exception:
        _RAPID_OCR_ENGINE = None
    return _RAPID_OCR_ENGINE


def _resolve_tesseract_binary() -> str:
    if os.path.isabs(OCR_TESSERACT_BIN) and os.path.exists(OCR_TESSERACT_BIN):
        return OCR_TESSERACT_BIN
    return shutil.which(OCR_TESSERACT_BIN) or OCR_TESSERACT_BIN


def _ocr_with_rapidocr(image) -> tuple:
    if np is None:
        return "", "rapidocr:missing_numpy"

    engine = _get_rapidocr_engine()
    if engine is None:
        return "", "rapidocr:unavailable"

    try:
        result, _ = engine(np.asarray(image))
        text, avg_score, line_count = _parse_rapidocr_result(result)
        if not text:
            return "", "rapidocr:no_text"
        return text, f"rapidocr|lines:{line_count}|avg_conf:{avg_score:.2f}"
    except Exception as exc:
        return "", f"rapidocr:error:{exc.__class__.__name__}"


def _ocr_with_tesseract(image) -> tuple:
    if pytesseract is None:
        return "", "tesseract:module_missing"

    binary = _resolve_tesseract_binary()
    if not binary or (not os.path.exists(binary) and not shutil.which(binary)):
        return "", "tesseract:binary_missing"

    try:
        pytesseract.pytesseract.tesseract_cmd = binary
        config = f"--oem {OCR_TESSERACT_OEM} --psm {OCR_TESSERACT_PSM}"
        text = pytesseract.image_to_string(
            image,
            lang=OCR_TESSERACT_LANG,
            config=config,
            timeout=OCR_TIMEOUT_SECONDS,
        )
        normalized = normalize_text(text)
        if not normalized:
            return "", "tesseract:no_text"
        return normalized, "tesseract"
    except Exception as exc:
        return "", f"tesseract:error:{exc.__class__.__name__}"


def _ocr_image_with_priority(image) -> tuple:
    traces: List[str] = []
    for engine in OCR_ENGINE_PRIORITY:
        if engine == "rapidocr":
            text, trace = _ocr_with_rapidocr(image)
        elif engine == "tesseract":
            text, trace = _ocr_with_tesseract(image)
        else:
            text, trace = "", f"{engine}:unsupported"

        traces.append(trace)
        if text:
            return text, ";".join(traces)

    return "", ";".join(traces) if traces else "ocr:no_engine"


def _merge_native_and_ocr_text(native_text: str, ocr_text: str) -> str:
    native = normalize_text(native_text)
    ocr = normalize_text(ocr_text)

    if not native:
        return ocr
    if not ocr:
        return native

    if ocr in native:
        return native
    if native in ocr:
        return ocr

    native_lines = {normalize_text(line) for line in native.splitlines() if normalize_text(line)}
    ocr_lines = [normalize_text(line) for line in ocr.splitlines() if normalize_text(line)]
    supplemental_lines = [line for line in ocr_lines if line not in native_lines]

    if not supplemental_lines:
        return native

    return f"{native}\n[OCR_SUPPLEMENT]\n" + "\n".join(supplemental_lines)


def _extract_ocr_text_from_pdf_page(pdfium_document: Any, page_index: int) -> tuple:
    if pdfium_document is None:
        return "", "ocr_pdf:pdfium_unavailable"
    if Image is None:
        return "", "ocr_pdf:pillow_unavailable"

    page = None
    bitmap = None
    try:
        page = pdfium_document[page_index]
        bitmap = page.render(scale=OCR_PDF_RENDER_DPI / 72.0)
        image = bitmap.to_pil()
        prepared = _prepare_image_for_ocr(image)
        return _ocr_image_with_priority(prepared)
    except Exception as exc:
        return "", f"ocr_pdf:error:{exc.__class__.__name__}"
    finally:
        try:
            if page is not None:
                page.close()
        except Exception:
            pass
        try:
            if bitmap is not None:
                bitmap.close()
        except Exception:
            pass


def _flatten_json(prefix: str, node: Any, out: List[str]) -> None:
    if isinstance(node, dict):
        for key, value in node.items():
            child_prefix = f"{prefix}.{key}" if prefix else str(key)
            _flatten_json(child_prefix, value, out)
        return

    if isinstance(node, list):
        for index, value in enumerate(node):
            child_prefix = f"{prefix}[{index}]"
            _flatten_json(child_prefix, value, out)
        return

    value = "" if node is None else str(node)
    normalized = normalize_text(value)
    if normalized:
        out.append(f"{prefix}: {normalized}" if prefix else normalized)


def _flatten_xml(root: etree._Element, out: List[str]) -> None:
    for element in root.iter():
        tag_name = str(element.tag)
        path = root.getroottree().getpath(element)
        if element.attrib:
            for key, value in element.attrib.items():
                normalized_value = normalize_text(str(value))
                if normalized_value:
                    out.append(f"{path}.@{key}: {normalized_value}")
        text_value = normalize_text(element.text or "")
        if text_value:
            out.append(f"{path}<{tag_name}>: {text_value}")


def extract_text_from_textlike(content: bytes, extension: str) -> Dict[str, Any]:
    ext = _normalize_extension(extension)
    strategy = _with_lib_versions(f"textlike:{ext}", "beautifulsoup4", "lxml", "pandas")
    if not content:
        return _response("", _error("ERR_EMPTY", "Empty"), strategy)

    decoded = safe_decode(content)
    if ext in {"html", "htm"}:
        decoded = strip_html_tags(decoded)
    elif ext == "rtf":
        decoded = _strip_rtf_control_words(decoded)
        strategy = f"{strategy}|rtf_stripped"
    elif ext == "json":
        try:
            parsed_json = json.loads(decoded)
            flattened_lines: List[str] = []
            _flatten_json("", parsed_json, flattened_lines)
            decoded = "\n".join(flattened_lines)
            strategy = f"{strategy}|flattened"
        except Exception as exc:
            strategy = f"{strategy}|json_fallback:{exc.__class__.__name__}"
    elif ext == "xml":
        try:
            root = etree.fromstring(decoded.encode("utf-8", errors="ignore"))
            flattened_lines: List[str] = []
            _flatten_xml(root, flattened_lines)
            decoded = "\n".join(flattened_lines)
            strategy = f"{strategy}|flattened"
        except Exception as exc:
            strategy = f"{strategy}|xml_fallback:{exc.__class__.__name__}"
    elif ext in {"csv", "tsv"}:
        try:
            dataframe = pd.read_csv(
                io.StringIO(decoded),
                sep="\t" if ext == "tsv" else ",",
                dtype=str,
                engine="python",
                quotechar='"',
                escapechar="\\",
                on_bad_lines="skip",
                keep_default_na=False,
                na_filter=False,
                nrows=MAX_CSV_ROWS,
            )

            lines = []
            headers = dataframe.columns.tolist()
            for row_values in dataframe.itertuples(index=False, name=None):
                row_context = ", ".join(
                    [
                        f"{str(header)}: {str(value)}"
                        for header, value in zip(headers, row_values)
                        if value is not None and str(value).strip() != ""
                    ]
                )
                if row_context:
                    lines.append(f"- {row_context}")

            if len(dataframe) >= MAX_CSV_ROWS:
                lines.append(f"[WARNING] CSV row limit reached ({MAX_CSV_ROWS} rows)")
            decoded = "\n".join(lines) if lines else decoded
        except Exception as exc:
            decoded_preview = _truncate_text(decoded)
            decoded = f"### CSV Data\n[TABLE_START]\n{decoded_preview}\n[TABLE_END]"
            strategy = f"{strategy}|csv_fallback:{exc.__class__.__name__}"

    decoded = normalize_text(decoded)
    return _response(decoded, None if decoded else _error("ERR_EMPTY", "Empty decoded"), strategy)


def extract_text_from_image_ocr(content: bytes, extension: str) -> Dict[str, Any]:
    ext = _normalize_extension(extension)
    strategy = _with_lib_versions(
        f"ocr_image:{ext}|engines:{','.join(OCR_ENGINE_PRIORITY)}",
        "pillow",
        "rapidocr-onnxruntime",
        "pytesseract",
    )

    if not content:
        return _response("", _error("ERR_EMPTY", "Empty"), strategy)
    if not OCR_ENABLED:
        return _response("", _error("ERR_OCR_DISABLED", "OCR disabled by configuration"), f"{strategy}|disabled")
    if Image is None:
        return _response("", _error("ERR_OCR_DEPENDENCY", "Pillow unavailable"), strategy)

    try:
        with Image.open(io.BytesIO(content)) as image_obj:
            prepared = _prepare_image_for_ocr(image_obj.copy())
    except Exception as exc:
        return _response("", _error("ERR_OCR_IMAGE", str(exc)), strategy)

    text, ocr_trace = _ocr_image_with_priority(prepared)
    if text:
        return _with_ocr_metadata(
            _response(text, None, f"{strategy}|{ocr_trace}"),
            attempted=True,
            used=True,
            trace=ocr_trace,
            pages=1,
            supplement_pages=1,
        )

    return _with_ocr_metadata(
        _response("", _error("ERR_OCR_EMPTY", "No text recognized"), f"{strategy}|{ocr_trace}"),
        attempted=True,
        used=False,
        trace=ocr_trace,
        pages=0,
        supplement_pages=0,
    )


def extract_text_from_pdf_pdfplumber(content: bytes, extension: str) -> Dict[str, Any]:
    ext = _normalize_extension(extension) or "pdf"
    strategy = _with_lib_versions(
        f"pdfplumber:{ext}|table_text_excluded:{int(PDF_EXCLUDE_TABLE_TEXT)}",
        "pdfplumber",
    )
    if not content:
        return _response("", _error("ERR_EMPTY", "Empty"), strategy)

    try:
        started_at = time.time()
        elements = []
        pdfium_document = None
        ocr_pages = 0
        ocr_supplement_pages = 0
        ocr_attempted_pages = 0
        ocr_traces: List[str] = []

        if OCR_ENABLED and OCR_PDF_FALLBACK_ENABLED and pdfium is not None:
            try:
                pdfium_document = pdfium.PdfDocument(content)
            except Exception as exc:
                strategy = f"{strategy}|ocr_pdfium_init:{exc.__class__.__name__}"

        with pdfplumber.open(io.BytesIO(content)) as pdf:
            total_pages = len(pdf.pages)
            page_limit_reached = total_pages > MAX_PDF_PAGES
            pages = pdf.pages[:MAX_PDF_PAGES]
            used_layout_fallback = False
            timeout_partial = False
            last_page_processed = 0

            if page_limit_reached:
                elements.append(f"[WARNING] PDF page limit reached ({MAX_PDF_PAGES}/{total_pages})")

            for page in pages:
                if _is_timed_out(started_at):
                    timeout_partial = True
                    elements.append(
                        f"[WARNING] PDF parsing timeout reached after {PARSE_TIMEOUT_SECONDS}s at page {last_page_processed}. Returning partial extraction."
                    )
                    break

                elements.append(f"## Page {page.page_number}")
                last_page_processed = page.page_number

                working_page = page
                try:
                    working_page = page.dedupe_chars(tolerance=1)
                except Exception:
                    pass

                table_settings = {
                    "vertical_strategy": "lines",
                    "horizontal_strategy": "lines",
                    "intersection_x_tolerance": 5,
                    "intersection_y_tolerance": 5,
                }

                try:
                    table_objects = working_page.find_tables(table_settings=table_settings)
                except Exception:
                    table_objects = []

                table_bboxes = [table.bbox for table in table_objects if getattr(table, "bbox", None)]

                page_text = ""
                if PDF_EXCLUDE_TABLE_TEXT and table_bboxes:
                    page_text = _extract_pdf_text_excluding_tables(working_page, table_bboxes)

                if not page_text:
                    page_text = working_page.extract_text(layout=False, x_tolerance=2, y_tolerance=2)
                if not page_text:
                    page_text = working_page.extract_text(layout=True, x_tolerance=2, y_tolerance=2)
                    if page_text:
                        used_layout_fallback = True

                native_char_count = len(normalize_text(page_text or ""))
                should_run_ocr = (
                    OCR_ENABLED
                    and OCR_PDF_FALLBACK_ENABLED
                    and page.page_number <= OCR_MAX_PDF_PAGES
                    and native_char_count < OCR_MIN_PDF_PAGE_TEXT_CHARS
                )

                if should_run_ocr:
                    ocr_attempted_pages += 1
                    ocr_text, ocr_trace = _extract_ocr_text_from_pdf_page(pdfium_document, page.page_number - 1)
                    ocr_traces.append(f"p{page.page_number}:{ocr_trace}")
                    if ocr_text:
                        ocr_pages += 1
                        if page_text:
                            merged_page_text = _merge_native_and_ocr_text(page_text, ocr_text)
                            if merged_page_text != page_text:
                                ocr_supplement_pages += 1
                            page_text = merged_page_text
                        else:
                            page_text = ocr_text
                            ocr_supplement_pages += 1

                if page_text:
                    elements.append(page_text)

                if table_objects:
                    extracted_tables = []
                    for table in table_objects:
                        try:
                            extracted_rows = table.extract()
                        except Exception:
                            extracted_rows = []
                        if extracted_rows:
                            extracted_tables.append(extracted_rows)
                else:
                    extracted_tables = working_page.extract_tables(table_settings=table_settings)

                for table in extracted_tables:
                    markdown_table = rows_to_markdown_table(table)
                    if markdown_table:
                        elements.append(f"\n[TABLE_START]\n{markdown_table}\n[TABLE_END]\n")

            if used_layout_fallback:
                strategy = f"{strategy}|layout_fallback"
            if page_limit_reached:
                strategy = f"{strategy}|page_limit:{MAX_PDF_PAGES}"
            if timeout_partial:
                strategy = f"{strategy}|timeout_partial:{last_page_processed}"
            if OCR_ENABLED and OCR_PDF_FALLBACK_ENABLED:
                strategy = f"{strategy}|ocr_pages:{ocr_pages}|ocr_supplement_pages:{ocr_supplement_pages}"
                if ocr_traces:
                    strategy = f"{strategy}|ocr_trace:{','.join(ocr_traces[:6])}"

        try:
            if pdfium_document is not None:
                pdfium_document.close()
        except Exception:
            pass

        full_text = normalize_text("\n\n".join(elements))
        ocr_trace_joined = ",".join(ocr_traces[:20]) if ocr_traces else None
        if not full_text and timeout_partial:
            return _with_ocr_metadata(
                _response("", _error("ERR_TIMEOUT", f"PDF parsing exceeded {PARSE_TIMEOUT_SECONDS}s"), strategy),
                attempted=ocr_attempted_pages > 0,
                used=ocr_pages > 0,
                trace=ocr_trace_joined,
                pages=ocr_pages,
                supplement_pages=ocr_supplement_pages,
            )
        return _with_ocr_metadata(
            _response(full_text, None if full_text else _error("ERR_EMPTY", "Empty PDF"), strategy),
            attempted=ocr_attempted_pages > 0,
            used=ocr_pages > 0,
            trace=ocr_trace_joined,
            pages=ocr_pages,
            supplement_pages=ocr_supplement_pages,
        )
    except Exception as exc:
        lowered_message = str(exc).lower()
        if "password" in lowered_message or "encrypted" in lowered_message:
            return _response("", _error("ERR_PDF_PASSWORD", str(exc)), strategy)
        return _response("", _error("ERR_PDF", str(exc)), strategy)


def extract_text_from_docx(content: bytes, extension: str) -> Dict[str, Any]:
    ext = _normalize_extension(extension) or "docx"
    strategy = _with_lib_versions(f"python_docx:{ext}", "python-docx")
    if not content:
        return _response("", _error("ERR_EMPTY", "Empty"), strategy)

    try:
        started_at = time.time()
        document = Document(io.BytesIO(content))
        elements = []
        ocr_attempted_units = 0
        ocr_used_units = 0
        ocr_traces: List[str] = []

        for section_index, section in enumerate(document.sections, start=1):
            header_text = normalize_text("\n".join([p.text for p in section.header.paragraphs if p.text.strip()]))
            footer_text = normalize_text("\n".join([p.text for p in section.footer.paragraphs if p.text.strip()]))
            if header_text:
                elements.append(f"## Header {section_index}\n{header_text}")
            if footer_text:
                elements.append(f"## Footer {section_index}\n{footer_text}")

        comments = _extract_docx_comments(document)
        if comments:
            elements.append("## Comments")
            elements.extend([f"- {comment}" for comment in comments])

        for child in document.element.body.iterchildren():
            if _is_timed_out(started_at):
                return _response("", _error("ERR_TIMEOUT", f"DOCX parsing exceeded {PARSE_TIMEOUT_SECONDS}s"), strategy)

            if child.tag.endswith("p"):
                from docx.text.paragraph import Paragraph

                paragraph = Paragraph(child, document)
                style_name = paragraph.style.name.lower() if paragraph.style else ""
                text = _extract_paragraph_with_hyperlinks(paragraph)
                if text:
                    if style_name.startswith("heading"):
                        level = 1
                        parts = style_name.split()
                        if len(parts) > 1 and parts[1].isdigit():
                            level = max(1, min(6, int(parts[1])))
                        text = f"{'#' * level} {text}"
                    elif "list" in style_name:
                        text = f"- {text}"
                    elements.append(text)
            elif child.tag.endswith("tbl"):
                from docx.table import Table

                rows = [[cell.text.strip() for cell in row.cells] for row in Table(child, document).rows]
                markdown_table = rows_to_markdown_table(rows)
                if markdown_table:
                    elements.append(f"\n[TABLE_START]\n{markdown_table}\n[TABLE_END]\n")

        textboxes = _extract_docx_textboxes(document)
        if textboxes:
            elements.append("## Textboxes")
            elements.extend([f"- {textbox}" for textbox in textboxes])

        image_blobs = _collect_docx_image_blobs(document)
        if OCR_ENABLED and image_blobs:
            strategy = f"{strategy}|docx_images:{len(image_blobs)}"
            ocr_attempted_units = len(image_blobs)

            if Image is None:
                ocr_traces.append("docx_images:pillow_unavailable")
                strategy = f"{strategy}|ocr_dependency_missing:pillow"
            else:
                ocr_sections = []
                for image_index, blob in enumerate(image_blobs, start=1):
                    try:
                        with Image.open(io.BytesIO(blob)) as image_obj:
                            prepared = _prepare_image_for_ocr(image_obj.copy())
                    except Exception as exc:
                        ocr_traces.append(f"img{image_index}:decode_error:{exc.__class__.__name__}")
                        continue

                    ocr_text, ocr_trace = _ocr_image_with_priority(prepared)
                    ocr_traces.append(f"img{image_index}:{ocr_trace}")
                    if not ocr_text:
                        continue

                    ocr_used_units += 1
                    ocr_sections.append(f"### Embedded image {image_index}\n{ocr_text}")

                if ocr_sections:
                    elements.append("## OCR embedded images")
                    elements.extend(ocr_sections)

            strategy = f"{strategy}|ocr_embedded_images:{ocr_used_units}/{len(image_blobs)}"
            if ocr_traces:
                strategy = f"{strategy}|ocr_trace:{','.join(ocr_traces[:6])}"

        full_text = normalize_text("\n\n".join(elements))
        return _with_ocr_metadata(
            _response(full_text, None if full_text else _error("ERR_EMPTY", "Empty Docx"), strategy),
            attempted=ocr_attempted_units > 0,
            used=ocr_used_units > 0,
            trace=",".join(ocr_traces[:20]) if ocr_traces else None,
            pages=ocr_used_units,
            supplement_pages=ocr_used_units,
        )
    except Exception as exc:
        return _response("", _error("ERR_DOCX", str(exc)), strategy)


def _resolve_antiword_binary() -> str:
    local_candidate = Path(__file__).resolve().parent.parent / ".tools" / "antiword" / "bin" / "antiword"
    if local_candidate.exists() and os.access(local_candidate, os.X_OK):
        return str(local_candidate)

    if os.path.isabs(ANTIWORD_BIN) and os.path.exists(ANTIWORD_BIN):
        return ANTIWORD_BIN
    which_result = shutil.which(ANTIWORD_BIN)
    return which_result or ANTIWORD_BIN


def _resolve_antiword_share_dir() -> str:
    if ANTIWORD_SHARE_DIR and os.path.isdir(ANTIWORD_SHARE_DIR):
        return ANTIWORD_SHARE_DIR

    local_share_candidate = Path(__file__).resolve().parent.parent / ".tools" / "antiword" / "share" / "antiword"
    if local_share_candidate.is_dir():
        return str(local_share_candidate)
    return ""


def extract_text_from_doc_antiword(content: bytes) -> Dict[str, Any]:
    strategy = f"antiword:doc|limits:cpu={ANTIWORD_CPU_SECONDS},mem={ANTIWORD_MEMORY_MB}MB"
    if not content:
        return _response("", _error("ERR_EMPTY", "Empty"), strategy)

    antiword_binary = _resolve_antiword_binary()
    antiword_share_dir = _resolve_antiword_share_dir()
    strategy = f"{strategy}|bin:{os.path.basename(str(antiword_binary))}"
    if not antiword_binary or (not shutil.which(antiword_binary) and not os.path.exists(antiword_binary)):
        return _response(
            "",
            _error("ERR_ANTIWORD_BIN", f"Antiword binary not found ({ANTIWORD_BIN}). Configure ANTIWORD_BIN environment variable."),
            strategy,
        )

    try:
        with tempfile.TemporaryDirectory(prefix="parsing_doc_") as temp_dir:
            temp_path = os.path.join(temp_dir, "input.doc")
            with open(temp_path, "wb") as temp_file:
                temp_file.write(content)

            custom_env = os.environ.copy()
            if antiword_share_dir:
                custom_env["ANTIWORDHOME"] = antiword_share_dir

            result = subprocess.run(
                [antiword_binary, "-m", "UTF-8.txt", "-w", "0", temp_path],
                capture_output=True,
                text=True,
                check=False,
                env=custom_env,
                timeout=PARSE_TIMEOUT_SECONDS,
                preexec_fn=_antiword_preexec_limits(),
            )

            if result.returncode == 0 and result.stdout.strip():
                return _response(normalize_text(result.stdout), None, strategy)

            error_text = result.stderr.strip() or f"Antiword failed with code {result.returncode}"
            return _response("", _error("ERR_ANTIWORD", error_text), strategy)
    except subprocess.TimeoutExpired:
        return _response("", _error("ERR_TIMEOUT", f"Antiword timed out after {PARSE_TIMEOUT_SECONDS}s"), strategy)
    except Exception as exc:
        return _response("", _error("ERR_ANTIWORD", str(exc)), strategy)


def extract_text_from_excel(content: bytes, extension: str) -> Dict[str, Any]:
    ext = _normalize_extension(extension) or "xlsx"
    engine = _excel_engine_for_extension(ext)
    excel_mode = "row_context+table" if EXCEL_INCLUDE_ROW_CONTEXT else "table_only"
    strategy = _with_lib_versions(f"pandas_excel:{ext}|engine:{engine or 'auto'}|mode:{excel_mode}", "pandas")
    if not content:
        return _response("", _error("ERR_EMPTY", "Empty"), strategy)

    try:
        started_at = time.time()

        try:
            workbook = pd.ExcelFile(io.BytesIO(content), engine=engine) if engine else pd.ExcelFile(io.BytesIO(content))
        except Exception:
            workbook = pd.ExcelFile(io.BytesIO(content))
            strategy = f"{strategy}|engine_fallback:auto"

        selected_sheets = workbook.sheet_names[:MAX_EXCEL_SHEETS]
        elements = []

        for sheet_name in selected_sheets:
            if _is_timed_out(started_at):
                return _response("", _error("ERR_TIMEOUT", f"Excel parsing exceeded {PARSE_TIMEOUT_SECONDS}s"), strategy)

            dataframe = workbook.parse(sheet_name=sheet_name, nrows=MAX_EXCEL_ROWS_PER_SHEET, dtype=str)
            elements.append(f"## Sheet: {sheet_name}")

            headers = dataframe.columns.tolist()
            preview_rows = [headers]

            for row_index, row_values in enumerate(dataframe.itertuples(index=False, name=None)):
                normalized_values = [normalize_text("" if value is None else str(value)) for value in row_values]
                if row_index < EXCEL_TABLE_PREVIEW_ROWS:
                    preview_rows.append(normalized_values)

                if EXCEL_INCLUDE_ROW_CONTEXT:
                    row_context = ", ".join(
                        [
                            f"{str(header)}: {str(value)}"
                            for header, value in zip(headers, normalized_values)
                            if value
                        ]
                    )
                    if row_context:
                        elements.append(f"- {row_context}")

            markdown_preview = rows_to_markdown_table(preview_rows)
            if markdown_preview:
                elements.append(f"\n[TABLE_START]\n{markdown_preview}\n[TABLE_END]\n")

            if len(dataframe) >= MAX_EXCEL_ROWS_PER_SHEET:
                elements.append(f"[WARNING] Row limit reached ({MAX_EXCEL_ROWS_PER_SHEET})")

        if len(workbook.sheet_names) > MAX_EXCEL_SHEETS:
            elements.append(f"[WARNING] Sheet limit reached ({MAX_EXCEL_SHEETS})")

        full_text = normalize_text("\n".join(elements))
        return _response(full_text, None if full_text else _error("ERR_EMPTY", "Empty Excel"), strategy)
    except Exception as exc:
        return _response("", _error("ERR_EXCEL", str(exc)), f"pandas_excel:{ext}")


def extract_text_from_pptx(content: bytes, extension: str) -> Dict[str, Any]:
    ext = _normalize_extension(extension) or "pptx"
    strategy = _with_lib_versions(f"python_pptx:{ext}", "python-pptx")
    if not content:
        return _response("", _error("ERR_EMPTY", "Empty"), strategy)

    try:
        started_at = time.time()
        presentation = Presentation(io.BytesIO(content))
        elements = []

        total_slides = len(presentation.slides)
        if total_slides > MAX_PPTX_SLIDES:
            elements.append(f"[WARNING] Slide limit reached ({MAX_PPTX_SLIDES}/{total_slides})")
            strategy = f"{strategy}|slide_limit:{MAX_PPTX_SLIDES}"

        def shape_sort_key(shape):
            return (
                _safe_int(getattr(shape, "top", 0)),
                _safe_int(getattr(shape, "left", 0)),
                _safe_int(getattr(shape, "shape_id", 0)),
            )

        def extract_chart_text(shape) -> List[str]:
            if not getattr(shape, "has_chart", False):
                return []

            chart = shape.chart
            chart_title = ""
            try:
                if chart.has_title and chart.chart_title and chart.chart_title.text_frame:
                    chart_title = normalize_text(chart.chart_title.text_frame.text)
            except Exception:
                chart_title = ""

            series_summaries = []
            for series in chart.series:
                series_name = normalize_text(str(getattr(series, "name", ""))) or "Series"
                try:
                    values = [str(value) for value in series.values]
                except Exception:
                    values = []

                if values:
                    preview = ", ".join(values[:12])
                    if len(values) > 12:
                        preview = f"{preview}, ..."
                    series_summaries.append(f"{series_name}: {preview}")

            if not series_summaries and chart_title:
                return [f"Chart: {chart_title}"]
            if series_summaries and chart_title:
                return [f"Chart: {chart_title} | {' ; '.join(series_summaries)}"]
            if series_summaries:
                return [f"Chart: {' ; '.join(series_summaries)}"]
            return ["Chart"]

        def extract_shape_texts(shape) -> Iterable[str]:
            if getattr(shape, "has_text_frame", False) and shape.text:
                text = normalize_text(shape.text)
                if text:
                    yield text

            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    rows.append([normalize_text(cell.text) for cell in row.cells])
                markdown_table = rows_to_markdown_table(rows)
                if markdown_table:
                    yield f"[TABLE_START]\n{markdown_table}\n[TABLE_END]"

            for chart_text in extract_chart_text(shape):
                normalized_chart = normalize_text(chart_text)
                if normalized_chart:
                    yield normalized_chart

            if hasattr(shape, "shapes"):
                for nested_shape in shape.shapes:
                    yield from extract_shape_texts(nested_shape)

        for slide_num, slide in enumerate(presentation.slides, start=1):
            if _is_timed_out(started_at):
                return _response("", _error("ERR_TIMEOUT", f"PPTX parsing exceeded {PARSE_TIMEOUT_SECONDS}s"), strategy)
            if slide_num > MAX_PPTX_SLIDES:
                break

            layout_name = normalize_text(getattr(slide.slide_layout, "name", ""))
            slide_title = ""
            try:
                if slide.shapes.title and slide.shapes.title.text:
                    slide_title = normalize_text(slide.shapes.title.text)
            except Exception:
                slide_title = ""

            slide_header = f"## Slide {slide_num}"
            if layout_name:
                slide_header = f"{slide_header} | Layout: {layout_name}"
            elements.append(slide_header)
            if slide_title:
                elements.append(f"### Slide Title: {slide_title}")

            for shape in sorted(list(slide.shapes), key=shape_sort_key):
                elements.extend([text for text in extract_shape_texts(shape) if text])

            if slide.has_notes_slide:
                notes_text = normalize_text(slide.notes_slide.notes_text_frame.text)
                if notes_text:
                    elements.append(f"## Slide {slide_num} Notes\n{notes_text}")

        full_text = normalize_text("\n\n".join(elements))
        return _response(full_text, None if full_text else _error("ERR_EMPTY", "Empty"), strategy)
    except Exception as exc:
        return _response("", _error("ERR_PPTX", str(exc)), strategy)


def extract_text_locally(content: bytes, extension: str) -> Dict[str, Any]:
    start_time = time.time()
    ext = _normalize_extension(extension)

    if not content:
        result = _response("", _error("ERR_EMPTY", "No content"), f"unhandled:{ext or 'none'}")
        result["parse_time_seconds"] = float(round(time.time() - start_time, 4))
        return result

    if len(content) > MAX_DOCUMENT_BYTES:
        result = _response(
            "",
            _error("ERR_SIZE_LIMIT", f"Document size {len(content)} exceeds limit {MAX_DOCUMENT_BYTES} bytes"),
            f"guard:size_limit:{ext or 'unknown'}",
        )
        result["parse_time_seconds"] = float(round(time.time() - start_time, 4))
        return result

    detected_format = _detect_binary_format(content)
    if _is_mime_mismatch_for_textlike(ext, detected_format):
        result = _response(
            "",
            _error("ERR_MIME_MISMATCH", f"Extension '{ext}' does not match binary signature '{detected_format}'"),
            f"guard:mime_mismatch:{ext}|signature:{detected_format}",
        )
        result["parse_time_seconds"] = float(round(time.time() - start_time, 4))
        return result

    extension_format_mismatch = (
        (ext == "pdf" and detected_format not in {"pdf", "unknown"})
        or (ext in {"doc", "xls", "ppt"} and detected_format not in {"ole2", "unknown"})
        or (ext in {"docx", "xlsx", "pptx", "xlsm", "xlsb"} and detected_format not in {"ooxml", "unknown"})
        or (ext in _IMAGE_EXTENSIONS and detected_format not in {"jpeg", "png", "tiff", "webp", "unknown"})
    )
    format_suffix = f"|signature:{detected_format}" if extension_format_mismatch else ""

    if ext == "txt":
        text = normalize_text(safe_decode(content))
        result = _response(text, None if text else _error("ERR_EMPTY", "Empty TXT"), f"txt_decoder{format_suffix}")
    elif ext == "pdf":
        result = extract_text_from_pdf_pdfplumber(content, ext)
        if extension_format_mismatch:
            result["parser_strategy"] = f"{result.get('parser_strategy', 'pdfplumber:pdf')}{format_suffix}"
    elif ext in {"docx", "docm"}:
        result = extract_text_from_docx(content, ext)
        if extension_format_mismatch:
            result["parser_strategy"] = f"{result.get('parser_strategy', 'python_docx:docx')}{format_suffix}"
    elif ext == "doc":
        result = extract_text_from_doc_antiword(content)
        if extension_format_mismatch:
            result["parser_strategy"] = f"{result.get('parser_strategy', 'antiword:doc')}{format_suffix}"
    elif ext in {"xls", "xlsx", "xlsm", "xlsb"}:
        result = extract_text_from_excel(content, ext)
        if extension_format_mismatch:
            result["parser_strategy"] = f"{result.get('parser_strategy', 'pandas_excel_row_context')}{format_suffix}"
    elif ext in {"pptx"}:
        result = extract_text_from_pptx(content, ext)
        if extension_format_mismatch:
            result["parser_strategy"] = f"{result.get('parser_strategy', 'python_pptx:pptx')}{format_suffix}"
    elif ext in _IMAGE_EXTENSIONS:
        result = extract_text_from_image_ocr(content, ext)
        if extension_format_mismatch:
            result["parser_strategy"] = f"{result.get('parser_strategy', 'ocr_image')}{format_suffix}"
    elif ext in TEXTLIKE_EXTENSIONS:
        result = extract_text_from_textlike(content, ext)
    elif detected_format in {"jpeg", "png", "tiff", "webp"}:
        result = extract_text_from_image_ocr(content, detected_format)
        result["parser_strategy"] = f"{result.get('parser_strategy', 'ocr_image')}|detected_format_routing"
    else:
        result = _response("", _error("ERR_UNHANDLED_EXTENSION", f"No local parser available for {ext}"), f"unhandled:{ext}")

    if len(result.get("text") or "") > MAX_CHUNK_CHARS * 200:
        result["parser_strategy"] = f"{result.get('parser_strategy', 'local')}|very_large_text"

    result["parse_time_seconds"] = float(round(time.time() - start_time, 4))
    return result


def extract_text_locally_with_retry(
    content: bytes,
    extension: str,
    max_retries: int = LOCAL_PARSE_MAX_RETRIES,
) -> Dict[str, Any]:
    started_at = time.time()
    attempts = max(1, int(max_retries))
    result: Dict[str, Any] = {}

    for attempt in range(1, attempts + 1):
        result = extract_text_locally(content, extension)
        parser_error = result.get("parser_error")
        if not parser_error:
            break
        if attempt >= attempts or not is_retryable_error(parser_error):
            break

    base_strategy = result.get("parser_strategy", "local")
    result["parser_strategy"] = f"{base_strategy}|attempts:{attempt}"
    result["parse_time_seconds"] = float(round(time.time() - started_at, 4))
    return result
